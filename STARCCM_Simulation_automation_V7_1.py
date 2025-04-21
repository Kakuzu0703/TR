import copy
import csv
import datetime
import json
import os
import shutil
import subprocess
import sys
import logging


from PyQt5.QtCore import Qt, QSize
from PyQt5.QtGui import QPixmap, QColor
from PyQt5.QtWidgets import QApplication, QWidget, QVBoxLayout, QHBoxLayout, QFormLayout, QLabel, QLineEdit, \
    QPushButton, QSizePolicy, QComboBox, QDialog, QMessageBox, QGroupBox, QListWidget, QListWidgetItem, QFileDialog
from ctREFPROP.ctREFPROP import REFPROPFunctionLibrary

from datetime import date
from pptx import Presentation
# import numpy



def change_unicode(han):
    unicode_str = ''.join(f'\\u{ord(char):04x}' if ord(char) > 127 else char for char in han)
    return unicode_str


# 安全计算表达式
def safe_eval(expr):
    try:
        return eval(expr)
    except:
        return None  # 或返回默认值/抛出异常

def get_formatted_date():
    today = date.today()
    return f"{today.year}.{today.month}.{today.day}"
def rename_and_save_step_file(source_path, destination_folder, new_name):
    """
    读取一个STEP文件，重命名后另存到指定文件夹

    :param source_path: 原始STEP文件的路径
    :param destination_folder: 目标文件夹路径
    :param new_name: 新的文件名（包括扩展名）
    """
    # # 确保目标文件夹存在
    # if not os.path.exists(destination_folder):
    #     os.makedirs(destination_folder)
    #     logging.info(f"文件夹已创建: {destination_folder}")
    # else:
    #     logging.info(f"文件夹已存在: {destination_folder}")

    # 构建目标文件的完整路径
    destination_path = os.path.join(destination_folder, new_name)

    # 检查源文件是否存在
    if not os.path.exists(source_path):
        logging.error(f"源文件不存在: {source_path}")
        return

        # 规范化路径并确保是字符串类型
    source_path = os.path.normpath(os.path.abspath(str(source_path)))
    destination_path = os.path.normpath(os.path.abspath(str(destination_path)))

    # 复制并重命名文件
    shutil.copy2(source_path, destination_path)
    logging.info(f"临时模型文件已保存完成: {destination_path}")

# 示例调用
# rename_and_save_step_file("C:/Users/owner/Desktop/star/改名模型/CV.STEP", "D:/STARCCM Simulation automation/CacheModels", "new_model_name.STEP")

def extract_model_name(file_path):
    # 获取文件名部分
    file_name = os.path.basename(file_path)
    # 去掉扩展名
    model_name = os.path.splitext(file_name)[0]
    return model_name
# # 示例调用
# file_path = r"C:\Users\owner\Desktop\star\改名模型\CV.STEP"
# model_name = extract_model_name(file_path)
# print(model_name)  # 输出: CV

def get_fluid_properties(RP_path, T_C, P_MPa,fluname):
    """
    通过REFPROP获取流体物性参数
    参数：
        RP_path: REFPROP DLL文件路径
        T_C: 温度（摄氏度）
        P_MPa: 压力（MPa）
    返回：
        density (kg/m³), viscosity (Pa·s)
    """
    # 初始化REFPROP
    RP = REFPROPFunctionLibrary(RP_path, 'dll')
    # 动态生成fluids目录路径
    fluid_dir = os.path.join(os.path.dirname(RP_path), "fluids")
    RP.SETPATHdll(fluid_dir)

    # 初始化工质（这里保持R134a，可根据需要参数化）
    r = RP.SETUPdll(1, f"{fluname}.FLD", "HMX.BNC", "DEF")
    if r.ierr != 0:
        raise ValueError(f"REFPROP初始化失败，错误代码：{r.ierr}")

    # 获取摩尔质量信息
    info = RP.INFOdll(1)

    # 单位转换
    T_K = T_C + 273.15
    P_kPa = P_MPa * 1000

    try:
        # 计算密度
        fla = RP.TPFLSHdll(T_K, P_kPa, [1.0])
        density_value = fla.D * info.wmm  # 转换为kg/m³
        density = float(f"{density_value:.5g}")

        # 计算粘度
        tr = RP.TRNPRPdll(T_K, fla.D, [1.0])
        viscosity_value = tr.eta * 1e-6  # 转换为Pa·s
        viscosity = float(f"{viscosity_value:.5g}")

        fla2 = RP.TPFLSHdll(T_K, P_kPa, [1.0])
        # 从返回结果获取声速（单位：m/s）
        speed_of_sound = float(f"{fla2.w:.5g}")  # w即为声速

        return density, viscosity,speed_of_sound

    except Exception as e:
        raise RuntimeError(f"物性计算失败: {str(e)}")
def read_last_row_last_column(csv_file_path):
    """
    读取CSV文件并返回最后一行最后一列的数据

    :param csv_file_path: CSV文件的路径
    :return: 最后一行最后一列的数据
    """
    try:
        with open(csv_file_path, mode='r', newline='', encoding='utf-8') as file:
            csv_reader = csv.reader(file)
            # 读取所有行
            rows = list(csv_reader)
            if not rows:
                logging.error(f"CSV文件为空: {csv_file_path}")
                return None
            # 提取最后一行最后一列的数据
            last_value = rows[-1][-1]
            return last_value
    except FileNotFoundError:
        logging.error(f"文件未找到: {csv_file_path}")
        return None
    except Exception as e:
        logging.error(f"读取CSV文件时发生错误: {e}")
        return None

def replace_image(slide, original_img_desc, new_img_path, target_index=None):
    """
    增强版图片替换函数（支持定位替换）
    :param original_img_desc: 要匹配的图片特征（支持部分文件名）
    :param new_img_path: 新图片完整路径
    :param target_index: 要替换的图片序号（从0开始，None表示替换全部）
    :return: 替换成功的图片数量
    """
    replaced_count = 0
    # 收集所有匹配图片
    targets = []
    for shape in slide.shapes:
        if shape.shape_type == 13 and original_img_desc in shape.image.filename:
            targets.append(shape)

    # 处理索引有效性
    if target_index is not None:
        if target_index >= len(targets) or target_index < 0:
            raise IndexError(f"无效索引：{target_index}，共找到{len(targets)}张匹配图片")
        targets = [targets[target_index]]  # 只保留指定索引的图片

    # 替换目标图片
    for shape in targets:
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height

        slide.shapes.add_picture(new_img_path, left, top, width, height)
        slide.shapes._spTree.remove(shape._element)
        replaced_count += 1

    return replaced_count

def modify_table(slide, row, col, new_text, table_index=0):
    """
        通用表格修改函数
        :param slide: 幻灯片对象
        :param row: 目标行号（从0开始）
        :param col: 目标列号（从0开始）
        :param new_text: 要更新的文本内容
        :param table_index: 表格索引（默认第1个表格）
        """
    """增强版表格修改（带颜色安全处理）"""
    tables = [shape for shape in slide.shapes if shape.has_table]
    if not tables:
        raise ValueError("幻灯片中未找到表格")
    table = tables[table_index].table

    # 行列索引修正（支持1-based索引）
    adj_row = row - 1
    adj_col = col - 1

    if adj_row >= len(table.rows) or adj_col >= len(table.columns):
        raise IndexError(f"无效坐标({row},{col})，表格尺寸{len(table.rows)}x{len(table.columns)}")

    cell = table.cell(adj_row, adj_col)

    # 清空单元格内容但保留格式
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = ""

    # 获取基准格式（使用第一个存在的run）
    base_run = None
    if cell.text_frame.paragraphs:
        paragraph = cell.text_frame.paragraphs[0]
        if paragraph.runs:
            base_run = paragraph.runs[0]
        else:
            base_run = paragraph.add_run()

    # 添加新内容
    new_run = cell.text_frame.paragraphs[0].add_run()
    new_run.text = new_text

    # 安全继承格式
    if base_run:
        new_run.font.bold = base_run.font.bold
        new_run.font.italic = base_run.font.italic
        new_run.font.size = base_run.font.size
        new_run.font.name = base_run.font.name

        # 颜色继承处理
        try:
            if base_run.font.color.rgb:
                new_run.font.color.rgb = base_run.font.color.rgb
            elif base_run.font.color.theme_color:
                new_run.font.color.theme_color = base_run.font.color.theme_color
            else:
                new_run.font.color.auto = True
        except AttributeError:
            new_run.font.color.auto = True


def append_text_to_slide(slide, target_text="压降仿真_", additional_text="（2024年最新数据）"):
    """
    安全颜色继承版本（支持所有颜色类型）
    """
    for shape in slide.shapes:
        if not (shape.has_text_frame and target_text in shape.text):
            continue

        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            if target_text not in paragraph.text:
                continue

            if paragraph.runs:
                last_run = paragraph.runs[-1]
                new_run = paragraph.add_run()
                new_run.text = additional_text

                # 字体基础属性继承
                new_run.font.bold = last_run.font.bold
                new_run.font.italic = last_run.font.italic
                new_run.font.size = last_run.font.size
                new_run.font.name = last_run.font.name

                # 安全颜色继承（关键修改部分）
                try:
                    # 优先继承RGB颜色
                    if last_run.font.color.rgb is not None:
                        new_run.font.color.rgb = last_run.font.color.rgb
                    # 其次继承主题颜色
                    elif last_run.font.color.theme_color is not None:
                        new_run.font.color.theme_color = last_run.font.color.theme_color
                    # 最后保持自动颜色
                    else:
                        new_run.font.color.auto = True
                except AttributeError as e:
                    # print(f"颜色继承异常: {str(e)}，已设为自动颜色")
                    new_run.font.color.auto = True
            else:
                new_run = paragraph.add_run()
                new_run.text = target_text + additional_text
            break
        break


# 在类定义顶部添加配置路径常量
CONFIG_FILE = "sim_config.json"

class SimulationConfigWindow(QWidget):
    def __init__(self, validator):#, validator):
        super().__init__()

        self.validator = validator  # 保存验证器实例

        self.last_input_params = {}  # 新增属性存储上次参数
        self.config = {}  # 确保已有配置字典

        self.task_queue = []  # 新增队列初始化

        # 新增路径变量初始化
        self.res_sce_path1 = None
        self.res_sce_path2 = None
        self.res_sce_path3 = None

        self.load_config()

        # 配置日志记录器（替换原来的basicConfig）
        self.logger = logging.getLogger()
        self.logger.setLevel(logging.INFO)

        # 移除已有处理器避免重复
        if self.logger.hasHandlers():
            self.logger.handlers.clear()

        # 创建控制台处理器
        console_handler = logging.StreamHandler()
        console_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))

        # # 创建文件处理器（按模型名生成日志文件）
        # log_path = f"D:\\STARCCM Simulation automation\\Log\\{name}.log"
        # file_handler = logging.FileHandler(log_path, encoding='utf-8')
        # file_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))

        # 添加处理器
        self.logger.addHandler(console_handler)
        # self.logger.addHandler(file_handler)

        # 初始化 index 变量
        self.index = 1

        # 初始化状态变量
        self.process_state = None

        # 原初始化代码替换为：
        self.model_import_path_input = QLineEdit()
        self.model_import_path_input.setText(self.config['model_import_path'])

        self.model_import_path_input.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Fixed)

        self.starccm_path_input = QLineEdit()
        self.starccm_path_input.setText(self.config['starccm_path'])

        self.starccmview_path_input= QLineEdit()
        self.starccmview_path_input.setText(self.config['starccmview_path'])

        self.refprp64dll_input = QLineEdit()
        self.refprp64dll_input.setText(self.config['refprp64dll'])

        self.threads_input = QLineEdit()
        self.threads_input.setText(self.config['threads'])

        self.stop_criteria_max_steps_input = QLineEdit()
        self.stop_criteria_max_steps_input.setText(self.config['max_steps'])

        self.workingfluid_input = QComboBox()
        self.workingfluid_input.addItems(["R134a", "R1234yf","R744", "50EG"])
        self.workingfluid_input.setCurrentIndex(self.config['workingfluid_index'])

        self.temperature_input = QLineEdit()
        self.temperature_input.setText(self.config['temperature'])

        self.pressure_input = QLineEdit()
        self.pressure_input.setText(self.config['pressure'])

        self.inlet_mass_flow_rate_input = QLineEdit()
        self.inlet_mass_flow_rate_input.setText(self.config['mass_flow'])

        self.operator_name_input= QLineEdit()
        self.operator_name_input.setText(self.config['operator_name'])

        # self.base_size_input = QLineEdit()
        # self.base_size_input.setText("1.0")
        #
        # self.target_surface_ratio_input = QLineEdit()
        # self.target_surface_ratio_input.setText("20.0")
        #
        # self.min_surface_ratio_input = QLineEdit()
        # self.min_surface_ratio_input.setText("8.0")
        #
        # self.prisma_layer_thickness_ratio_input = QLineEdit()
        # self.prisma_layer_thickness_ratio_input.setText("6.0")
        #
        # self.prisma_layer_extension_input = QLineEdit()
        # self.prisma_layer_extension_input.setText("2")

        # self.viscosity_input = QLineEdit()
        # self.viscosity_input.setText("1.1276E-5")
        #
        # self.density_input = QLineEdit()
        # self.density_input.setText("19.077")


        # 初始化图片标签
        self.pressure_label = QLabel()
        self.pressure_label.setFixedSize(300, 200)
        self.streamline_label = QLabel()
        self.streamline_label.setFixedSize(300, 200)

        # self.STATUS_COLOR = {
        #     "等待计算": "#95a5a6",  # 灰色
        #     "计算中": "#3498db",  # 蓝色
        #     "已完成": "#2ecc71",  # 绿色
        #     "失败": "#e74c3c"  # 红色
        # }
        #
        # self.STATUS_BG_COLOR = {
        #     "等待计算": "#f0f0f0",  # 浅灰
        #     "计算中": "#e3f2fd",  # 淡蓝
        #     "已完成": "#e8f5e9",  # 淡绿
        #     "失败": "#ffebee"  # 淡红
        # }
        self.STATUS_COLOR = {
            "等待计算": "#ffffff",  # 白色
            "计算中": "#ffffff",   # 白色
            "已完成": "#ffffff",   # 白色
            "失败": "#ffffff"     # 白色
        }

        self.STATUS_BG_COLOR = {
            "等待计算": "#FFA500",  # 橙色
            "计算中": "#3498db",   # 蓝色
            "已完成": "#2ecc71",   # 绿色
            "失败": "#e74c3c"     # 红色
        }

        self.initUI()
        self.load_config()




    def load_config(self):
        # 完整默认配置
        default_config = {
            'model_import_path': r'C:\Users\owner\Desktop\star\改名模型\CV.STEP',
            'starccm_path': r'D:\starCCM+\16.06.008-R8\STAR-CCM+16.06.008-R8\star\lib\win64\intel20.1vc14.2-r8\lib\starccmw.exe',
            'starccmview_path': r'D:\starCCM+\16.06.008-R8\STAR-View+16.06.008\bin\starview+.exe',
            'refprp64dll': r'D:\Program Files (x86)\REFPROP\REFPRP64.DLL',
            'threads': "64",
            'max_steps': "3000",
            'temperature': "25",
            'pressure': "0.3",
            'mass_flow': "210/3600",
            'workingfluid_index': 0,
            'operator_name': "",
            'last_params': {},  # 新增参数存储
            'task_queue': []
        }

        try:
            if os.path.exists(CONFIG_FILE):
                with open(CONFIG_FILE, 'r', encoding='utf-8') as f:
                    saved_config = json.load(f)
                    # 合并配置（保证默认配置完整性）
                    self.config = {**default_config, **saved_config}
                    self.last_input_params = saved_config.get('last_params', {})
                    self.task_queue = saved_config.get('task_queue', [])
                # 在加载配置后添加队列初始化逻辑
                if hasattr(self, 'queue_list'):  # 确保UI组件已初始化
                    self.queue_list.clear()
                    for task in self.task_queue:
                        if 'submit_time' in task:
                            task['submit_time'] = datetime.datetime.strptime(
                                task['submit_time'], "%Y-%m-%d %H:%M:%S"
                            )
                        # 重构任务信息生成逻辑
                        task_info = (
                            f"任务 {self.task_queue.index(task) + 1}: "
                            f"{extract_model_name(task['model_import_path'])} | "
                            f"操作员：{task['operator_name']} | "
                            f"工质: {['R134a', 'R1234yf', 'R744', '50EG'][task['workingfluid_index']]} | "  # 通过索引获取工质名称
                            f"温度: {task['temperature']}℃ | "
                            f"绝对压力: {task['pressure']}MPa | "
                            f"质量流量: {task['mass_flow']}kg/s | "
                            f"提交时间: {task['submit_time']}"
                        )
                        # 使用HTML实现状态颜色
                        # colored_text = f'<font color="{self.STATUS_COLOR[task["status"]]}">[{task["status"]}]</font>'
                        colored_text = (
                            f'<span style="'
                            f'background-color: {self.STATUS_BG_COLOR[task["status"]]}; '
                            f'color: {self.STATUS_COLOR[task["status"]]}; '
                            f'padding: 2px 5px; border-radius: 3px;'
                            f'">[{task["status"]}]</span>'
                        )
                        full_text = f"{colored_text} {task_info}"

                        item = QListWidgetItem()

                        # 创建自定义Widget实现部分着色
                        label = QLabel(full_text)
                        label.setTextFormat(Qt.RichText)
                        item.setSizeHint(QSize(label.sizeHint().width(), label.sizeHint().height() + 10))
                        self.queue_list.addItem(item)
                        self.queue_list.setItemWidget(item, label)
                        self.update_queue_display()
            else:
                self.config = default_config
                self.last_input_params = {}
                self.task_queue = []
        except Exception as e:
            print(f"加载配置失败: {str(e)}")
            self.config = default_config
            self.last_input_params = {}
            self.task_queue = []

    def save_config(self):
        try:
            save_data = {
                # 保留所有配置字段
                'model_import_path': self.model_import_path_input.text().strip('"'),
                'starccm_path': self.starccm_path_input.text().strip('"'),
                'starccmview_path': self.starccmview_path_input.text().strip('"'),
                'refprp64dll': self.refprp64dll_input.text().strip('"'),
                'threads': self.threads_input.text(),
                'max_steps': self.stop_criteria_max_steps_input.text(),
                'temperature': self.temperature_input.text(),
                'pressure': self.pressure_input.text(),
                'mass_flow': self.inlet_mass_flow_rate_input.text(),
                'workingfluid_index': self.workingfluid_input.currentIndex(),
                'operator_name': self.operator_name_input.text(),
                'last_params': self.last_input_params,  # 新增参数存储
                'task_queue': [
            {k: v.strftime("%Y-%m-%d %H:%M:%S") if isinstance(v, datetime.datetime) else v
             for k, v in task.items()}  # 时间转字符串
            for task in self.task_queue
        ]
            }
            with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
                json.dump(save_data, f, indent=2, ensure_ascii=False)
        except Exception as e:
            print(f"保存配置失败: {str(e)}")

    def closeEvent(self, event):
        """窗口关闭时自动保存"""
        self.save_config()
        event.accept()

    # 在SimulationConfigWindow类中添加新方法：
    def open_scene_file(self, scene_path):
        """使用STAR-View+打开场景文件"""
        try:
            if os.path.exists(scene_path):
                viewer_path = self.starccmview_path_input.text().strip('"')
                subprocess.Popen([viewer_path, scene_path])
            else:
                QMessageBox.warning(self, "错误", f"场景文件不存在:\n{scene_path}")
        except Exception as e:
            logging.error(f"打开场景文件失败: {str(e)}")

    def get_scene_path(self, filename,params):
        name = extract_model_name(params['model_import_path'])
        operator_name = params['operator_name']
        datenow = get_formatted_date()
        return os.path.join(
            "D:", "仿真自动化结果",
            datenow, operator_name,
            f"{name}_{self.index}",  # 这里会动态获取最新的 self.index
            "Report",
            f"{name}_{filename}"
        )

    def get_pressure_drop_display(self, Ma, last_value):
        if Ma == 0:
            return f"平均压降值: {int(round(float(last_value)))} Pa"
        elif Ma == 1:
            return f"平均压降值: {int(round(float(last_value)))} Pa\n警告！内部部分流速马赫数大于0.3！"
        elif Ma == 2:
            return "平均压降值: N/A\n警告！内部部分流速马赫数大于0.5！"
        else:
            return "平均压降值: 数据异常"

    def add_to_queue(self):

        #先保存至最后一次输入
        self.last_input_params={
            'model_import_path': self.model_import_path_input.text().strip('"'),
            'starccm_path': self.starccm_path_input.text().strip('"'),
            'starccmview_path': self.starccmview_path_input.text().strip('"'),
            'refprp64dll': self.refprp64dll_input.text().strip('"'),
            'threads': self.threads_input.text(),
            'max_steps': self.stop_criteria_max_steps_input.text(),
            'temperature': self.temperature_input.text(),
            'pressure': self.pressure_input.text(),
            'mass_flow': self.inlet_mass_flow_rate_input.text(),
            'workingfluid': self.workingfluid_input.currentText(),
            'operator_name': self.operator_name_input.text()
        }
        # self.save_config()

        """捕获当前输入参数并添加到队列"""
        current_params1 = {
                'status': "等待计算",
                'model_import_path': self.model_import_path_input.text().strip('"'),
                'starccm_path': self.starccm_path_input.text().strip('"'),
                'starccmview_path': self.starccmview_path_input.text().strip('"'),
                'refprp64dll': self.refprp64dll_input.text().strip('"'),
                'threads': self.threads_input.text(),
                'max_steps': self.stop_criteria_max_steps_input.text(),
                'temperature': self.temperature_input.text(),
                'pressure': self.pressure_input.text(),
                'mass_flow': self.inlet_mass_flow_rate_input.text(),
                'workingfluid_index': self.workingfluid_input.currentIndex(),
                'operator_name': self.operator_name_input.text(),
                'submit_time': datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),  # 新增提交时间
                'simulation_index': None,
                'simulation_date': None,
                'Ma': None,
                'res_mach_number': None
        }

        # ==== 新增校验逻辑 ====
        if not self.operator_name_input.text().strip():
            QMessageBox.warning(self, "输入错误", "请先输入操作员姓名")
            return

        # 修改后的重复任务检查（忽略status字段）
        temp_params = current_params1.copy()
        temp_params.pop('status', None)
        temp_params.pop('submit_time', None)
        temp_params.pop('simulation_index', None)
        temp_params.pop('simulation_date', None)
        temp_params.pop('Ma', None)
        temp_params.pop('res_mach_number', None)

        for idx, task in enumerate(self.task_queue, 1):
            task_copy = task.copy()
            task_copy.pop('status', None)
            task_copy.pop('submit_time', None)
            task_copy.pop('simulation_index', None)
            task_copy.pop('simulation_date', None)
            task_copy.pop('Ma', None)
            task_copy.pop('res_mach_number', None)
            if task_copy == temp_params:
                reply = QMessageBox.question(
                    self, "重复任务",
                    f"任务 {idx} 与当前参数完全相同\n\n是否继续添加？",
                    QMessageBox.Yes | QMessageBox.No,
                    QMessageBox.No
                )
                if reply == QMessageBox.No:
                    return

        # 深度拷贝避免参数被后续修改
        self.task_queue.append(copy.deepcopy(current_params1))

        # 显示在队列列表
        task_info = (f"任务 {len(self.task_queue)}: "
                     f'''{extract_model_name(self.model_import_path_input.text().strip('"'))} | '''
                     f"操作员：{current_params1['operator_name']} | "
                     f"工质: {self.workingfluid_input.currentText()} | "
                     f"温度: {current_params1['temperature']}℃ | "
                     f"绝对压力: {current_params1['pressure']}MPa | "
                     f"质量流量: {current_params1['mass_flow']}kg/s | "
                     f"提交时间: {current_params1['submit_time']}")
        # 使用HTML实现部分文字着色
        status_bg = self.STATUS_BG_COLOR[current_params1["status"]]
        status_color = self.STATUS_COLOR[current_params1["status"]]
        colored_text = (
            f'<span style="background-color: {status_bg}; color: {status_color}; '
            f'padding: 2px 5px; border-radius: 3px;">[{current_params1["status"]}]</span>'
        )
        full_text = f"{colored_text} {task_info}"

        item = QListWidgetItem()
        item.setData(Qt.UserRole, current_params1)

        # 创建自定义Widget实现部分着色
        label = QLabel(full_text)
        label.setTextFormat(Qt.RichText)
        item.setSizeHint(QSize(label.sizeHint().width(), label.sizeHint().height() + 10))

        self.queue_list.addItem(item)
        self.queue_list.setItemWidget(item, label)
        self.save_config()

    def delete_from_queue(self):
        # 检查队列是否为空
        if not self.task_queue:
            QMessageBox.warning(self, "提示", "任务队列已为空")
            return

        # 获取当前选中项（支持未选中时删除第一个）
        selected_index = self.queue_list.currentRow()
        if selected_index == -1:  # 当未选中任何项时
            selected_index = 0  # 默认删除第一个

        # 有效性验证
        if 0 <= selected_index < len(self.task_queue):
            # 同时删除数据队列和列表项
            del self.task_queue[selected_index]
            self.queue_list.takeItem(selected_index)
            self.save_config()
        else:
            QMessageBox.warning(self, "错误", "无效的任务索引")

        # 删除后刷新队列序号
        self.update_queue_display()

    def update_queue_display(self):
        """刷新队列显示序号"""
        self.queue_list.clear()
        for index, task in enumerate(self.task_queue, 1):
            task_info = (
                f"任务 {index}: "
                f"{extract_model_name(task['model_import_path'])} | "
                f"操作员：{task['operator_name']} | "
                f"工质: {['R134a', 'R1234yf', 'R744', '50EG'][task['workingfluid_index']]} | "
                f"温度: {task['temperature']}℃ | "
                f"绝对压力: {task['pressure']}MPa | "
                f"质量流量: {task['mass_flow']}kg/s | "
                f"提交时间: {task['submit_time']}"
            )

            # 新增着色逻辑
            colored_text = (
                f'<span style="'
                f'background-color: {self.STATUS_BG_COLOR[task["status"]]}; '
                f'color: {self.STATUS_COLOR[task["status"]]}; '
                f'padding: 2px 5px; border-radius: 3px;'
                f'">[{task["status"]}]</span>'
            )
            full_text = f"{colored_text} {task_info}"

            # 创建新item和label
            item = QListWidgetItem()
            label = QLabel(full_text)
            label.setTextFormat(Qt.RichText)
            item.setSizeHint(QSize(label.sizeHint().width(), label.sizeHint().height() + 10))

            self.queue_list.addItem(item)
            self.queue_list.setItemWidget(item, label)

    def clear_queue(self):
        self.task_queue.clear()
        self.queue_list.clear()
        self.update_queue_display()
        self.save_config()

    def select_model_file(self):
        options = QFileDialog.Options()
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            "选择模型文件",
            "",
            "STEP Files (*.step *.stp);;All Files (*)",
            options=options
        )
        if file_path:
            # 保持路径格式统一
            normalized_path = os.path.normpath(file_path)
            self.model_import_path_input.setText(normalized_path)

    def on_task_double_clicked(self, item):
        """双击任务项处理"""
        selected_index = self.queue_list.currentRow()

        if 0 <= selected_index < len(self.task_queue):
            task = self.task_queue[selected_index]

            # 如果任务已完成，直接显示结果
            # if task['status'] == "已完成":
                # self.load_task_params(selected_index)
                # self.display_simulation_results(task)
                # return

            # 未完成的任务保持原逻辑
            reply = QMessageBox.question(
                self, "加载任务参数",
                "是否要查看此任务的输入条件及结果？\n选择\"是\"将加载参数到输入框",
                QMessageBox.Yes | QMessageBox.No,
                QMessageBox.No
            )
            if reply == QMessageBox.Yes:
                self.load_task_params(selected_index)
                if task['status'] == "已完成":
                    self.display_simulation_results(task)

    def display_simulation_results(self, task):
        self.btn_mach_3d.setVisible(False)
    #     """显示仿真结果"""
        try:
            # 获取模型基本信息
            res_name = extract_model_name(task['model_import_path'])
            res_operator_name = task['operator_name']
            res_datenow = task['simulation_date']
            res_index = task['simulation_index']
            res_Ma=task['Ma']
            res_mach_number=task['res_mach_number']
            res_workflow = task['workingfluid_index']


            # 构建结果路径
            res_report_folder = os.path.join(
                "D:\\仿真自动化结果",
                res_datenow,
                res_operator_name,
                f"{res_name}_{res_index}",
                "Report"
            )

            # 显示压力云图
            res_pressure_img = os.path.join(res_report_folder, f"{res_name}_压力云图.png")
            if os.path.exists(res_pressure_img):
                pixmap = QPixmap(res_pressure_img).scaled(400, 250, Qt.KeepAspectRatio)
                self.pressure_label.setPixmap(pixmap)
                self.pressure_label.mouseDoubleClickEvent = lambda e: self.show_image(res_pressure_img)

            # 显示流线图
            res_streamline_img = os.path.join(res_report_folder, f"{res_name}_流线图.png")
            if os.path.exists(res_streamline_img):
                pixmap = QPixmap(res_streamline_img).scaled(400, 250, Qt.KeepAspectRatio)
                self.streamline_label.setPixmap(pixmap)
                self.streamline_label.mouseDoubleClickEvent = lambda e: self.show_image(res_streamline_img)

            # 显示压降值
            res_csv_path = os.path.join(res_report_folder, f"{res_name}_average_pressure.csv")
            # res_last_value = round(float(read_last_row_last_column(res_csv_path)))
            res_last_value = read_last_row_last_column(res_csv_path)
            if res_last_value is not None:
                res_last_value = round(float(read_last_row_last_column(res_csv_path)))
                self.pressure_drop_label.setText(f" {res_last_value} Pa")
                self.pressure_drop_label.setStyleSheet("font-size: 24px; color: #009900; font-weight: bold;")
            else:
                self.pressure_drop_label.setText("N/A 马赫数>0.5")
                self.pressure_drop_label.setStyleSheet("font-size: 24px; color: #FF0000; font-weight: bold;")

            res_color = "#000000"  # 默认黑色
            if res_workflow in [0,1,2]:
                if res_Ma==0:
                    res_color = "#009900"  # 绿色
                elif res_Ma == 1:
                    self.btn_mach_3d.setVisible(True)
                    self.btn_mach_3d.setEnabled(True)
                    res_color = "#FFA500"
                else:
                    self.btn_mach_3d.setVisible(True)
                    self.btn_mach_3d.setEnabled(True)
                    res_color = "#FF0000"

                self.mach_number_label.setText(f"{res_mach_number}")
                self.mach_number_label.setStyleSheet(f"""
                    font-size: 24px; 
                    color: {res_color}; 
                    font-weight: bold;
                """)
            else:
                self.mach_number_label.setText(f"{res_mach_number}")
                self.mach_number_label.setStyleSheet("font-size: 24px; color: #666666; font-weight: bold;")

            self.res_sce_path1 = os.path.join(res_report_folder, f"{res_name}_压力云图.sce")
            self.res_sce_path2 = os.path.join(res_report_folder, f"{res_name}_流线图.sce")
            self.res_sce_path3 = os.path.join(res_report_folder, f"{res_name}_Ma_0.3区域图.sce")


            self.btn_pressure_3d.setEnabled(True)
            self.btn_streamline_3d.setEnabled(True)
            QApplication.processEvents()  # 强制刷新UI


        except Exception as e:
            QMessageBox.warning(self, "结果加载失败",
                                f"无法加载仿真结果：{str(e)}\n"
                                f"可能原因：\n"
                                "1. 结果文件未正常生成\n"
                                "2. 文件路径被修改\n"
                                "3. 仿真未完全完成")

    def load_task_params(self, index):
        """加载指定任务的参数到输入框"""
        task = self.task_queue[index]

        try:
            # 通用参数设置
            self.model_import_path_input.setText(task['model_import_path'])
            self.starccm_path_input.setText(task['starccm_path'])
            self.starccmview_path_input.setText(task['starccmview_path'])
            self.refprp64dll_input.setText(task['refprp64dll'])
            self.threads_input.setText(task['threads'])
            self.stop_criteria_max_steps_input.setText(task['max_steps'])

            # 工质选择需要处理索引
            self.workingfluid_input.setCurrentIndex(task['workingfluid_index'])

            # 温度压力等数值参数
            self.temperature_input.setText(task['temperature'])
            self.pressure_input.setText(task['pressure'])
            self.inlet_mass_flow_rate_input.setText(task['mass_flow'])

            # 操作员姓名
            self.operator_name_input.setText(task['operator_name'])

            QMessageBox.information(self, "加载成功", "任务参数及结果已加载到输入框")
        except Exception as e:
            QMessageBox.critical(self, "加载错误", f"参数加载失败：{str(e)}")

    def initUI(self):
        # 设置窗口大小
        self.resize(1600, 900)  # 设置窗口宽度为1600像素，高度为400像素

        # 初始化 index 变量
        # self.index = 1

        # 创建背景图片标签
        background_label = QLabel(self)
        background_pixmap = QPixmap(
            os.path.join(os.path.dirname(__file__), 'logo_SANHUA.png'),
            format='PNG',
            flags=Qt.ImageConversionFlag.ColorOnly
        ).scaled(75, 75, Qt.KeepAspectRatio,
                                                                     Qt.SmoothTransformation)
        background_label.setPixmap(background_pixmap)
        background_label.setGeometry(0, 0, 75, 75)  # 设置图片的位置和大小

        # 创建标题标签
        title_label = QLabel('STAR-CCM+阀类压降仿真自动化')
        title_label.setStyleSheet("font-size: 36px; font-weight: bold; text-align: center;")

        # 主垂直布局
        main_layout = QVBoxLayout(self)
        main_layout.setContentsMargins(20, 20, 20, 20)
        main_layout.setSpacing(20)

        top_widget = QWidget()
        top_layout = QHBoxLayout(top_widget)
        top_layout.setContentsMargins(0, 10, 0, 0)
        top_layout.setSpacing(30)

        # 创建左侧部分
        left_layout = QVBoxLayout()
        left_layout.setSpacing(20)  # 设置垂直间距
        left_layout.setContentsMargins(20, 20, 20, 20)  # 设置边距：左, 上, 右, 下

        left_form_layout = QFormLayout()
        left_form_layout.setSpacing(20)  # 设置表单间距
        left_form_layout.setContentsMargins(20, 20, 20, 20)  # 设置边距：左, 上, 右, 下
        # left_form_layout.addRow(QLabel('导入数模路径:'), self.model_import_path_input)
        # 改为：
        file_selector_layout = QHBoxLayout()
        btn_choose_file = QPushButton("...")
        btn_choose_file.setStyleSheet("font-size: 20px; padding: 5px;")
        btn_choose_file.clicked.connect(self.select_model_file)  # 连接点击事件
        file_selector_layout.addWidget(self.model_import_path_input)
        file_selector_layout.addWidget(btn_choose_file)
        left_form_layout.addRow(QLabel('导入数模路径:'), file_selector_layout)

        left_form_layout.addRow(QLabel('STAR-CCM软件路径:'), self.starccm_path_input)
        left_form_layout.addRow(QLabel('STAR-CCM Viewer软件路径:'), self.starccmview_path_input)
        left_form_layout.addRow(QLabel('REFPRP64.DLL路径:'), self.refprp64dll_input)
        left_form_layout.addRow(QLabel('线程数:'), self.threads_input)
        left_form_layout.addRow(QLabel('停止准则 最大步数:'), self.stop_criteria_max_steps_input)

        # 物性参数设置
        left_form_layout2 = QFormLayout()
        left_form_layout2.setSpacing(20)  # 设置表单间距
        left_form_layout2.setContentsMargins(20, 20, 20, 20)  # 设置边距：左, 上, 右, 下
        property_title_label = QLabel('流体工质：')
        property_title_label.setStyleSheet("font-size: 28px; font-weight: bold;")
        left_form_layout2.addRow(property_title_label)
        # left_form_layout2.addRow(QLabel('动力粘度（Pa·s）:'), self.viscosity_input)
        # left_form_layout2.addRow(QLabel('密度（kg/m³）:'), self.density_input)
        left_form_layout2.addRow(QLabel('工质选择:'), self.workingfluid_input)

        left_layout.addLayout(left_form_layout)
        left_layout.addLayout(left_form_layout2)

        # 创建右侧部分
        right_layout = QVBoxLayout()
        right_layout.setSpacing(20)  # 设置垂直间距
        right_layout.setContentsMargins(20, 20, 20, 20)  # 设置边距：左, 上, 右, 下

        # 网格设置
        right_form_layout = QFormLayout()
        right_form_layout.setSpacing(20)  # 设置表单间距
        right_form_layout.setContentsMargins(20, 20, 20, 20)  # 设置边距：左, 上, 右, 下
        grid_title_label = QLabel('工况输入：')
        grid_title_label.setStyleSheet("font-size: 28px; font-weight: bold;")
        right_form_layout.addRow(grid_title_label)
        # right_form_layout.addRow(QLabel('基础尺寸（mm）:'), self.base_size_input)
        # right_form_layout.addRow(QLabel('目标表面尺寸 基数百分比:'), self.target_surface_ratio_input)
        # right_form_layout.addRow(QLabel('最小表面尺寸 基数百分比:'), self.min_surface_ratio_input)
        # right_form_layout.addRow(QLabel('棱柱层总厚度 基数百分比:'), self.prisma_layer_thickness_ratio_input)
        # right_form_layout.addRow(QLabel('棱柱层 层数:'), self.prisma_layer_extension_input)
        right_form_layout.addRow(QLabel('温度（℃）:'), self.temperature_input)
        right_form_layout.addRow(QLabel('入口绝对压力（MPa）:'), self.pressure_input)
        right_form_layout.addRow(QLabel('入口质量流量（kg/s）:'), self.inlet_mass_flow_rate_input)

        right_form_layout2 = QFormLayout()
        right_form_layout2.setSpacing(20)  # 设置表单间距
        right_form_layout2.setContentsMargins(20, 20, 20, 20)  # 设置边距：左, 上, 右, 下

        # grid1_title_label = QLabel('操作员：')
        # grid1_title_label.setStyleSheet("font-size: 28px; font-weight: bold;")
        # right_form_layout2.addRow(grid1_title_label)

        # 操作员姓名输入
        # self.operator_name_input = QLineEdit()
        self.operator_name_input.setPlaceholderText("请输入姓名")
        # self.operator_name_input.setStyleSheet("font-size: 20px;")
        operator_label = QLabel('操作员:')
        operator_label.setStyleSheet("font-size: 28px; font-weight: bold;")  # 从20px放大到28px，加粗
        right_form_layout2.addRow(operator_label, self.operator_name_input)

        # 在right_form_layout之后添加压降显示区域
        # right_form_layout.addRow(QLabel('当前压降值(Pa):'), QLabel())  # 先添加一个空Label占位

        # # 修改后的压降值布局
        # pressure_row = QHBoxLayout()
        # pressure_row.addWidget(QLabel('压降值(Pa):'))
        # self.pressure_drop_label = QLabel("  等待计算...")
        # self.pressure_drop_label.setStyleSheet("font-size: 24px; color: #FF0000; font-weight: bold;")
        # pressure_row.addWidget(self.pressure_drop_label)
        # pressure_row.addStretch()  # 添加弹性空间使内容左对齐
        #
        # # 将水平布局添加到表单布局
        # right_form_layout.addRow(pressure_row)

        # # 初始条件和边界条件设置
        # right_form_layout2 = QFormLayout()
        # right_form_layout2.setSpacing(20)  # 设置表单间距
        # right_form_layout2.setContentsMargins(20, 20, 20, 20)  # 设置边距：左, 上, 右, 下
        # initial_conditions_title_label = QLabel('初始条件和边界条件设置：')
        # initial_conditions_title_label.setStyleSheet("font-size: 28px; font-weight: bold;")
        # right_form_layout2.addRow(initial_conditions_title_label)
        # right_form_layout2.addRow(QLabel('入口压力（MPa）:'), self.pressure_input)
        # right_form_layout2.addRow(QLabel('入口质量流量（kg/s）:'), self.inlet_mass_flow_rate_input)

        # # 图片显示区域
        # image_container = QWidget()
        # image_layout = QHBoxLayout(image_container)
        # image_layout.setSpacing(50)  # 增加间距到50像素
        #
        # # 压力云图标签
        # self.pressure_label = QLabel("压力云图预览区")
        # self.pressure_label.setStyleSheet("font-size: 25px; font-weight: bold;")  # 增大字体
        # self.pressure_label.setAlignment(Qt.AlignCenter)  # 文字居中
        #
        # # 流线图标签
        # self.streamline_label = QLabel("流线图预览区")
        # self.streamline_label.setStyleSheet("font-size: 25px; font-weight: bold;")  # 增大字体
        # self.streamline_label.setAlignment(Qt.AlignCenter)  # 文字居中
        # image_layout.addWidget(self.pressure_label)
        # image_layout.addWidget(self.streamline_label)


        right_layout.addLayout(right_form_layout)
        right_layout.addLayout(right_form_layout2)
        # right_layout.addWidget(image_container)
        # right_layout.addLayout(right_form_layout2)

        # # 在类属性中添加结果标签
        # self.pressure_drop_label = QLabel("等待计算...")
        # self.pressure_drop_label.setStyleSheet("font-size: 24px; color: #FF0000; font-weight: bold;")
        # right_layout.insertWidget(2, self.pressure_drop_label)  # 插入到工况输入下方

        # 设置所有其他标签的字体大小为20号
        for label in [self.model_import_path_input, self.starccm_path_input,self.refprp64dll_input,self.starccmview_path_input, self.threads_input,
                      self.stop_criteria_max_steps_input, self.workingfluid_input,
                      self.temperature_input, self.pressure_input, self.inlet_mass_flow_rate_input,self.operator_name_input]:
            label.setStyleSheet("font-size: 20px;")

        # 创建开始运行按钮
        self.run_button = QPushButton('开始运行')
        self.run_button.setFixedSize(500, 50)  # 设置按钮的固定大小
        self.run_button.setStyleSheet("font-size: 48px; font-weight: bold;background-color: #90EE90;")  # 设置按钮的字体大小和背景颜色
        self.run_button.clicked.connect(self.on_button_clicked)  # 连接按钮点击事件

        # 创建一个空的QWidget来占据两行的空间
        spacer = QWidget()
        spacer.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)

        # 创建一个垂直布局来容纳表单布局和按钮
        right_bottom_layout = QVBoxLayout()
        right_bottom_layout.setSpacing(20)  # 设置垂直间距
        right_bottom_layout.addWidget(spacer)  # 添加一个占位符
        right_bottom_layout.addWidget(self.run_button)  # 添加按钮

        # 将表单布局和按钮布局添加到右侧布局中
        right_layout.addLayout(right_bottom_layout)

        # ==================== 队列管理布局 ====================
        queue_group = QGroupBox()#("仿真队列管理")
        # queue_group.setStyleSheet("""
        #     font-size: 28px;
        #     font-weight: bold;
        #     color: #2c3e50;
        #     padding-bottom: 10px;
        # """)
        queue_layout = QVBoxLayout()
        queue_layout.setContentsMargins(20, 10, 20, 10)  # 统一设置边距
        queue_layout.setSpacing(20)

        # 标题标签
        queue_title_label = QLabel('仿真队列管理：')
        queue_title_label.setStyleSheet("""
            font-size: 28px; 
            font-weight: bold;
            color: #2c3e50;
            padding-bottom: 10px;
        """)
        queue_layout.addWidget(queue_title_label)

        # 队列列表
        self.queue_list = QListWidget()
        self.queue_list.setFixedHeight(300)
        self.queue_list.setStyleSheet("""
            QListWidget {{
                border: 1px solid #bdc3c7;
                border-radius: 5px;
                background-color: #ffffff;
            }}
        """)
        queue_layout.addWidget(self.queue_list)

        # 按钮容器
        button_container = QWidget()
        button_layout = QHBoxLayout(button_container)
        button_layout.setContentsMargins(0, 0, 0, 0)
        button_layout.setSpacing(10)

        btn_add = QPushButton("＋ 添加任务")
        btn_delete = QPushButton("－ 删除任务")
        btn_clear = QPushButton("× 清空队列")

        # 按钮样式
        button_style = """
            QPushButton {{
                font-size: 14px;
                padding: 8px 15px;
                border-radius: 4px;
                min-width: 100px;
            }}
            QPushButton:hover {{
                opacity: 0.9;
            }}
        """
        btn_add.setStyleSheet(button_style + "background-color: #3498db; color: black;")
        btn_delete.setStyleSheet(button_style + "background-color: #e67e22; color: black;")
        btn_clear.setStyleSheet(button_style + "background-color: #e74c3c; color: black;")

        btn_add.clicked.connect(self.add_to_queue)
        btn_delete.clicked.connect(self.delete_from_queue)
        btn_clear.clicked.connect(self.clear_queue)
        self.queue_list.itemDoubleClicked.connect(self.on_task_double_clicked)

        button_layout.addWidget(btn_add)
        button_layout.addWidget(btn_delete)
        button_layout.addWidget(btn_clear)
        queue_layout.addWidget(button_container)

        queue_group.setLayout(queue_layout)

        # 创建主布局
        # main_layout = QHBoxLayout()
        # main_layout.setSpacing(30)  # 设置水平间距
        # main_layout.setContentsMargins(30, 30, 30, 30)  # 设置边距：左, 上, 右, 下
        # main_layout.addLayout(left_layout)
        # main_layout.addLayout(right_layout)

        top_layout.addLayout(left_layout)
        top_layout.addLayout(right_layout)
        top_layout.insertWidget(2, queue_group)

        # ==================== 下半部分 - 结果展示区 ====================
        bottom_widget = QWidget()
        bottom_layout = QVBoxLayout(bottom_widget)
        bottom_layout.setContentsMargins(20, 20, 20, 20)

        # 添加标题
        result_title = QLabel('输出：')
        result_title.setStyleSheet("""
            font-size: 28px; 
            font-weight: bold;
            margin-bottom: 15px;
        """)
        bottom_layout.addWidget(result_title)

        # 压降值显示
        pressure_drop_container = QWidget()
        drop_layout = QHBoxLayout(pressure_drop_container)
        drop_layout.addWidget(QLabel('压降值(Pa):'))
        self.pressure_drop_label = QLabel("等待计算...")
        self.pressure_drop_label.setStyleSheet("""
                font-size: 24px; 
                color: #FF0000; 
                font-weight: bold;
                padding: 5px 15px;
                background: #FFF3CD;
                border-radius: 5px;
            """)
        drop_layout.addWidget(self.pressure_drop_label)
        drop_layout.addStretch()

        # 马赫数组件
        drop_layout.addWidget(QLabel('马赫数:'))
        self.mach_number_label = QLabel("等待计算...")
        self.mach_number_label.setStyleSheet("""
                font-size: 24px; 
                color: #FF0000; 
                font-weight: bold;
                padding: 5px 15px;
                background: #FFF3CD;
                border-radius: 5px;
        """)
        drop_layout.addWidget(self.mach_number_label)
        drop_layout.addStretch()

        # 图片预览区
        image_container = QWidget()
        image_layout = QHBoxLayout(image_container)
        image_layout.setContentsMargins(0, 20, 0, 0)

        # 压力云图
        self.pressure_label = QLabel("压力云图预览区")
        self.pressure_label.setStyleSheet("""
               QLabel {
                   border: 2px solid #3498DB;
                   border-radius: 5px;
                   background: #F8F9FA;
                   min-width: 400px;
                   min-height: 250px;
                   font-size: 36px;
                   font-weight: bold;
               }
           """)
        self.pressure_label.setAlignment(Qt.AlignCenter)

        # 流线图
        self.streamline_label = QLabel("流线图预览区")
        self.streamline_label.setStyleSheet("""
               QLabel {
                   border: 2px solid #27AE60;
                   border-radius: 5px;
                   background: #F8F9FA;
                   min-width: 400px;
                   min-height: 250px;
                   font-size: 36px;
                   font-weight: bold;
               }
           """)
        self.streamline_label.setAlignment(Qt.AlignCenter)

        image_layout.addWidget(self.pressure_label)
        image_layout.addWidget(self.streamline_label)

        # 组合底部布局
        bottom_layout.addWidget(pressure_drop_container)
        bottom_layout.addWidget(image_container)

        # ==================== 3D可视化按钮布局 ====================
        visualization_widget = QWidget()
        visualization_layout = QHBoxLayout(visualization_widget)
        visualization_layout.setContentsMargins(20, 20, 20, 20)
        visualization_layout.setSpacing(50)

        # 添加标题
        tdresult_title = QLabel('三维结果可视化：')
        tdresult_title.setStyleSheet("""
            font-size: 28px; 
            font-weight: bold;
            margin-bottom: 15px;
        """)
        visualization_layout.addWidget(tdresult_title)

        # 创建按钮
        self.btn_pressure_3d = QPushButton('三维压力云图')
        self.btn_pressure_3d.setFixedSize(250, 50)  # 宽度280像素，高度60像素
        self.btn_streamline_3d = QPushButton('三维流线图')
        self.btn_streamline_3d.setFixedSize(250, 50)
        self.btn_mach_3d = QPushButton('马赫数>0.3区域')
        self.btn_mach_3d.setFixedSize(250, 50)

        # 设置按钮样式
        button_style = """
            QPushButton {
                font-size: 24px;
                font-weight: bold;
                padding: 15px 30px;
                border-radius: 8px;
                background-color: #3498DB;
                color: white;
            }
            QPushButton:hover {
                background-color: #2980B9;
            }
        """
        self.btn_pressure_3d.setStyleSheet(button_style)
        self.btn_streamline_3d.setStyleSheet(button_style)
        self.btn_mach_3d.setStyleSheet(button_style)

        # name = extract_model_name(self.model_import_path_input.text().strip('"'))
        # 读取操作员名字
        # operator_name = self.operator_name_input.text().strip()
        # if not operator_name:
        #     QMessageBox.warning(self, "输入错误", "请先输入操作员姓名")
        #     return
        datenow = get_formatted_date()


        # 设置按钮点击事件
        self.btn_pressure_3d.clicked.connect(
            # lambda: self.open_scene_file(self.get_scene_path("压力云图.sce")))
            lambda: self.open_scene_file(self.res_sce_path1)
        )

        self.btn_streamline_3d.clicked.connect(
            # lambda: self.open_scene_file(self.get_scene_path("流线图.sce")))
            lambda: self.open_scene_file(self.res_sce_path2)
        )

        self.btn_mach_3d.clicked.connect(
            # lambda: self.open_scene_file(self.get_scene_path("Ma_0.3区域图.sce")))
            lambda: self.open_scene_file(self.res_sce_path3)
        )

        # 初始隐藏马赫数按钮
        self.btn_mach_3d.setVisible(False)

        # 添加按钮到布局
        visualization_layout.addWidget(self.btn_pressure_3d)
        visualization_layout.addWidget(self.btn_streamline_3d)
        visualization_layout.addWidget(self.btn_mach_3d)
        self.btn_pressure_3d.setEnabled(False)
        self.btn_streamline_3d.setEnabled(False)
        self.btn_mach_3d.setEnabled(False)
        visualization_layout.addStretch()


        # ==================== 组合主布局 ====================
        main_layout.addWidget(top_widget, stretch=1)
        main_layout.addWidget(bottom_widget, stretch=1)
        main_layout.addWidget(visualization_widget, stretch=1)

        self.setLayout(main_layout)
        self.setWindowTitle('STAR-CCM+ 阀类压降仿真自动化V1.7')
        self.show()

    def show_image(self, image_path):
        """显示大图的对话框"""
        dialog = QDialog(self)
        # 从路径获取文件名作为标题
        dialog.setWindowTitle(os.path.basename(image_path))
        layout = QVBoxLayout(dialog)

        pixmap = QPixmap(image_path)
        if pixmap.isNull():
            QMessageBox.warning(self, "错误", "图片加载失败")
            return

        label = QLabel()
        # label.setPixmap(pixmap.scaled(1600, 900, aspectRatioMode=Qt.KeepAspectRatio))
        # 使用原始分辨率显示
        label.setPixmap(pixmap)  # 移除scaled缩放
        # 根据图片尺寸调整对话框
        dialog.resize(pixmap.width(), pixmap.height())
        layout.addWidget(label)

        dialog.exec_()

    def on_button_clicked(self):
        # pending_task = next((task for task in self.task_queue
        #                      if task.status == "等待计算"), None)
        # # 检查任务队列中是否存在等待计算的任务
        # self.operator_name_input.clear()
        pending_tasks = [task for task in self.task_queue if task.get('status') == "等待计算"]

        if pending_tasks:
            # 弹出确认对话框
            reply = QMessageBox.question(
                self, "任务队列检测",
                f"检测到 {len(pending_tasks)} 个待计算任务，是否开始自动执行队列计算？",
                QMessageBox.Yes | QMessageBox.No,
                QMessageBox.Yes
            )
            is_queue_task=True

            if reply == QMessageBox.Yes:
                while True:
                    # 查找第一个等待计算的任务
                    pending_tasks = [t for t in self.task_queue if t['status'] == "等待计算"]
                    if not pending_tasks:
                        break
                    # 取第一个等待任务
                    task = pending_tasks[0]
                    params = {
                        'model_import_path': task['model_import_path'],
                        'starccm_path': task['starccm_path'],
                        'starccmview_path': task['starccmview_path'],
                        'refprp64dll': task['refprp64dll'],
                        'threads': task['threads'],
                        'max_steps': task['max_steps'],
                        'temperature': task['temperature'],
                        'pressure': task['pressure'],
                        'mass_flow': task['mass_flow'],
                        'workingfluid_index': task['workingfluid_index'],
                        'operator_name': task['operator_name']
                    }
                    try:
                        # 更新任务状态为计算中
                        task['status'] = "计算中"
                        self.update_queue_display()
                        QApplication.processEvents()  # 强制UI刷新

                        # 执行仿真（需要将原有执行逻辑封装成独立方法）
                        self.on_run_button_clicked(params,is_queue_task=True)

                        if self.process_state == 1:
                            # 标记任务完成
                            task['status'] = "已完成"
                    except Exception as e:
                        logging.error(f"任务执行失败: {str(e)}")
                        task['status'] = "失败"

                    finally:
                        self.update_queue_display()
                        self.save_config()
                        QApplication.processEvents()
        else:
            # 弹出确认对话框
            reply = QMessageBox.question(
                self, "任务队列检测",
                "仿真队列内未检测到待计算任务，是否以当前输入条件进行仿真？",
                QMessageBox.Yes | QMessageBox.No,
                QMessageBox.Yes
            )
            is_queue_task=False

            if reply == QMessageBox.Yes:
                # 从输入框获取参数
                params = {
                    'model_import_path': self.model_import_path_input.text().strip('"'),
                    'starccm_path': self.starccm_path_input.text().strip('"'),
                    'starccmview_path': self.starccmview_path_input.text().strip('"'),
                    'refprp64dll': self.refprp64dll_input.text().strip('"'),
                    'threads': self.threads_input.text(),
                    'max_steps': self.stop_criteria_max_steps_input.text(),
                    'temperature': self.temperature_input.text(),
                    'pressure': self.pressure_input.text(),
                    'mass_flow': self.inlet_mass_flow_rate_input.text(),
                    'workingfluid_index': self.workingfluid_input.currentIndex(),
                    'operator_name': self.operator_name_input.text()
                }
                self.on_run_button_clicked(params,is_queue_task=False)
        # return params
    def on_run_button_clicked(self,params,is_queue_task):

        # 移除现有的日志处理器
        for handler in self.logger.handlers[:]:
            self.logger.removeHandler(handler)
            handler.close()
        console_handler = logging.StreamHandler()
        console_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
        # 添加处理器
        self.logger.addHandler(console_handler)

        start_time = datetime.datetime.now()  # 添加在方法开始处

        # ==== 新增校验逻辑 ====
        if not self.operator_name_input.text().strip():
            QMessageBox.warning(self, "输入错误", "请先输入操作员姓名")
            return
        # ==== 校验结束 ====

        # # ==== 新增参数对比逻辑 ====
        # current_params = {
        #     'model_path': self.model_import_path_input.text().strip('"'),
        #     'threads': self.threads_input.text(),
        #     'max_steps': self.stop_criteria_max_steps_input.text(),
        #     'temperature': self.temperature_input.text(),
        #     'pressure': self.pressure_input.text(),
        #     'mass_flow': self.inlet_mass_flow_rate_input.text(),
        #     'workingfluid': self.workingfluid_input.currentText(),
        #     'operator': self.operator_name_input.text()
        #     # 'starccm_path': self.starccm_path_input.text().strip('"'),
        #     # 'starccmview_path': self.starccmview_path_input.text().strip('"'),
        #     # 'refprp64dll_path': self.refprp64dll_input.text().strip('"'),
        # }
        #
        # if current_params == self.last_input_params and self.last_input_params:
        #     reply = QMessageBox.question(
        #         self, '输入确认',
        #         "输入参数与上次运行完全相同，是否继续？",
        #         QMessageBox.Yes | QMessageBox.No,
        #         QMessageBox.No
        #     )
        #     if reply == QMessageBox.No:
        #         return

        # ==== 新增参数对比逻辑 ====
        current_params = {
            'model_import_path': self.model_import_path_input.text().strip('"'),
            'starccm_path': self.starccm_path_input.text().strip('"'),
            'starccmview_path': self.starccmview_path_input.text().strip('"'),
            'refprp64dll': self.refprp64dll_input.text().strip('"'),
            'threads': self.threads_input.text(),
            'max_steps': self.stop_criteria_max_steps_input.text(),
            'temperature': self.temperature_input.text(),
            'pressure': self.pressure_input.text(),
            'mass_flow': self.inlet_mass_flow_rate_input.text(),
            'workingfluid': self.workingfluid_input.currentText(),
            'operator_name': self.operator_name_input.text()
        }

        # 转换为数值类型进行比较（处理字符串与数值的对比）
        try:
            current_params['temperature'] = float(current_params['temperature'])
            current_params['pressure'] = float(current_params['pressure'])
            current_params['mass_flow'] = float(current_params['mass_flow'])
        except ValueError:
            pass
        try:
            self.last_input_params['temperature'] = float(self.last_input_params['temperature'])
            self.last_input_params['pressure'] = float(self.last_input_params['pressure'])
            self.last_input_params['mass_flow'] = float(self.last_input_params['mass_flow'])
        except ValueError:
            pass

        if current_params == self.last_input_params and self.last_input_params and not is_queue_task:
            reply = QMessageBox.question(
                self,
                '确认输入',
                "输入参数与上次运行完全相同，是否继续？\n\n",
                # f"模型路径: {current_params['model_import_path']}\n"
                # f"工质类型: {current_params['workingfluid']}\n"
                # f"温度: {current_params['temperature']}℃",
                QMessageBox.Yes | QMessageBox.No,
                QMessageBox.No
            )
            if reply == QMessageBox.No:
                return

        # 初始隐藏马赫数按钮
        self.btn_mach_3d.setVisible(False)
        self.btn_pressure_3d.setEnabled(False)
        self.btn_streamline_3d.setEnabled(False)
        self.btn_mach_3d.setEnabled(False)

        # 重置压降值显示
        self.pressure_drop_label.setText(" 计算中...")
        self.pressure_drop_label.setStyleSheet("""
                font-size: 24px; 
                color: #FF0000; 
                font-weight: bold;
                padding: 5px 15px;
                background: #FFF3CD;
                border-radius: 5px;
            """)
        # 重置马赫数显示
        self.mach_number_label.setText(" 计算中...")
        self.mach_number_label.setStyleSheet("""
                font-size: 24px; 
                color: #FF0000; 
                font-weight: bold;
                padding: 5px 15px;
                background: #FFF3CD;
                border-radius: 5px;
        """)

        # 重置图片预览区
        self.pressure_label.clear()
        self.pressure_label.setText("压力云图预览区")
        self.pressure_label.setStyleSheet("""
               QLabel {
                   border: 2px solid #3498DB;
                   border-radius: 5px;
                   background: #F8F9FA;
                   min-width: 400px;
                   min-height: 250px;
                   font-size: 36px;
                   font-weight: bold;
               }
           """)

        self.streamline_label.clear()
        self.streamline_label.setText("流线图预览区")
        self.streamline_label.setStyleSheet("""
               QLabel {
                   border: 2px solid #27AE60;
                   border-radius: 5px;
                   background: #F8F9FA;
                   min-width: 400px;
                   min-height: 250px;
                   font-size: 36px;
                   font-weight: bold;
               }
           """)

        # 移除鼠标事件绑定
        self.pressure_label.mouseDoubleClickEvent = None
        self.streamline_label.mouseDoubleClickEvent = None

        # 立即刷新界面
        QApplication.processEvents()

        # 修改按钮文本和样式
        self.run_button.setText('运行中请勿点击')
        self.run_button.setStyleSheet("font-size: 48px; font-weight: bold; background-color: #FFA07A;")  # 浅红色背景
        self.run_button.setEnabled(False)  # 禁用按钮，防止多次点击
        QApplication.processEvents()  # 强制刷新UI

        # D盘创建一个仿真文件夹
        folder_path = "D:\\STARCCM Simulation automation"#加密文件夹

        # 检查并创建主文件夹
        if not os.path.exists(folder_path):
            os.makedirs(folder_path)
            logging.info(f"文件夹已创建: {folder_path}")
        else:
            logging.info(f"文件夹已存在: {folder_path}")

            # 获取当前日期并创建日期文件夹
        date_folder = os.path.join(folder_path, get_formatted_date())
        if not os.path.exists(date_folder):
            os.makedirs(date_folder)
            logging.info(f"日期文件夹已创建: {date_folder}")
        else:
            logging.info(f"日期文件夹已存在: {date_folder}")

        # 读取操作员名字
        operator_name = params['operator_name']
        if not operator_name:
            QMessageBox.warning(self, "输入错误", "请先输入操作员姓名")
            return

        # 创建操作员文件夹
        operator_folder = os.path.join(date_folder, operator_name)
        if not os.path.exists(operator_folder):
            os.makedirs(operator_folder)
            logging.info(f"操作员文件夹已创建: {operator_folder}")
        else:
            logging.info(f"操作员文件夹已存在: {operator_folder}")


        ###创建公开文件夹
        public_folder_path = "D:\\仿真自动化结果"

        # 检查并创建主文件夹
        if not os.path.exists(public_folder_path):
            os.makedirs(public_folder_path)
            logging.info(f"文件夹已创建: {public_folder_path}")
        else:
            logging.info(f"文件夹已存在: {public_folder_path}")

            # 获取当前日期并创建日期文件夹
        date_folder_public = os.path.join(public_folder_path, get_formatted_date())
        if not os.path.exists(date_folder_public):
            os.makedirs(date_folder_public)
            logging.info(f"日期文件夹已创建: {date_folder_public}")
        else:
            logging.info(f"日期文件夹已存在: {date_folder_public}")

        # 创建操作员文件夹
        operator_folder_public = os.path.join(date_folder_public, operator_name)
        if not os.path.exists(operator_folder_public):
            os.makedirs(operator_folder_public)
            logging.info(f"操作员文件夹已创建: {operator_folder_public}")
        else:
            logging.info(f"操作员文件夹已存在: {operator_folder_public}")


        name = extract_model_name(params['model_import_path'])
        # 创建模型文件夹
        # model_folder = os.path.join(operator_folder, name)
        current_index = 0
        while True:
            current_index += 1  # 先递增索引
            model_folder = os.path.join(operator_folder, f"{name}_{current_index}")
            model_folder_public = os.path.join(operator_folder_public, f"{name}_{current_index}")
            if not os.path.exists(model_folder):
                break
        self.index = current_index  # 保持index与文件夹一致

        # 创建模型文件夹
        os.makedirs(model_folder)
        logging.info(f"模型文件夹已创建: {model_folder}")

        os.makedirs(model_folder_public)
        logging.info(f"模型文件夹已创建: {model_folder_public}")

        # 定义子文件夹路径
        models_folder = os.path.join(model_folder, "CacheModels")
        simulation_folder = os.path.join(model_folder, "Simulation")
        report_folder = os.path.join(model_folder, "Report")
        log_folder = os.path.join(model_folder, "Log")

        # 检查并创建 Models 子文件夹
        if not os.path.exists(models_folder):
            os.makedirs(models_folder)
            logging.info(f"子文件夹已创建: {models_folder}")

        # 检查并创建 Simulation 子文件夹
        if not os.path.exists(simulation_folder):
            os.makedirs(simulation_folder)
            logging.info(f"子文件夹已创建: {simulation_folder}")

        # 检查并创建 Report 子文件夹
        if not os.path.exists(report_folder):
            os.makedirs(report_folder)
            logging.info(f"子文件夹已创建: {report_folder}")

        # 检查并创建 Log 子文件夹
        if not os.path.exists(log_folder):
            os.makedirs(log_folder)
            logging.info(f"子文件夹已创建: {log_folder}")


        # 定义子文件夹路径(公开)
        report_folder_public = os.path.join(model_folder_public, "Report")
        log_folder_public = os.path.join(model_folder_public, "Log")

        # 检查并创建 Report 子文件夹
        if not os.path.exists(report_folder_public):
            os.makedirs(report_folder_public)
            logging.info(f"子文件夹已创建: {report_folder_public}")

        # 检查并创建 Log 子文件夹
        if not os.path.exists(log_folder_public):
            os.makedirs(log_folder_public)
            logging.info(f"子文件夹已创建: {log_folder_public}")

        # 读取每个 QLineEdit 的值
        # model_import_path = self.model_import_path_input.text().strip('"')#.replace("\\", "/")
        # starccm_path = self.starccm_path_input.text().strip('"')#.replace("\\", "/")
        # refprop_path = self.refprp64dll_input.text().strip('"')#.replace("\\", "/"))
        # threads = self.threads_input.text()
        # stop_criteria_max_steps = self.stop_criteria_max_steps_input.text()
        model_import_path=params['model_import_path']
        starccm_path=params['starccm_path']
        refprop_path=params['refprp64dll']
        threads=params['threads']
        stop_criteria_max_steps=params['max_steps']

        base_size = 0.3#self.base_size_input.text()
        target_surface_ratio = 100#self.target_surface_ratio_input.text()
        min_surface_ratio =25# self.min_surface_ratio_input.text()
        prisma_layer_thickness_ratio = 33#self.prisma_layer_thickness_ratio_input.text()
        prisma_layer_extension = 2#self.prisma_layer_extension_input.text()

        # viscosity = self.viscosity_input.text()
        # density = self.density_input.text()
        temperature = float(params['temperature'])
        pressure = float(params['pressure'])
        workingfluid_index = params['workingfluid_index']
        # 在类中定义工质列表作为常量（推荐放在类顶部）
        WORKING_FLUIDS = ["R134a", "R1234yf", "R744", "50EG"]
        # 在需要获取工质名称的地方使用索引获取
        workingfluid = WORKING_FLUIDS[workingfluid_index]
        inlet_mass_flow_rate = params['mass_flow']

        workingfluid_1 = 'CO2' if workingfluid == 'R744' else workingfluid

        # 添加工质判断逻辑
        if workingfluid in ["R134a", "R1234yf","R744"]:
            dp_unit='"bar"'
            dp_format='"%-6.2f"'
            try:
                # 调用物性计算函数
                density, viscosity, speed_of_sound = get_fluid_properties(
                    RP_path=refprop_path,
                    T_C=temperature,
                    P_MPa=pressure,
                    fluname=workingfluid_1
                )
                logging.info(f"成功获取 {workingfluid} 物性参数：密度={density} kg/m³，粘度={viscosity} Pa·s, 声速={speed_of_sound} m/s")
            except Exception as e:
                logging.error(f"物性计算失败: {str(e)}")
                return
        else:
            dp_unit='"Pa"'
            dp_format='"%-6.0f"'
            speed_of_sound = 897
            # # 对于其他工质（如50EG），使用默认值或单独的处理逻辑
            # viscosityf = -0.000000000000002891086856001190294218*temperature ** 7+0.000000000001077136958500754662705 * temperature ** 6 - 0.000000000151185766170254471016 * temperature ** 5 + 0.00000000991277733151393456151 * temperature ** 4 - 0.0000003214977807661548076655 * temperature ** 3 + 0.000007525808299778965107627 * temperature ** 2 - 0.0002892859986891702268885 * temperature + 0.008247019621432158734131  # 示例值，需要根据实际情况修改
            # densityf = -0.00000000373756542382253*temperature**3-0.00243168387588982000000*temperature**2-0.33814727958631500000000*temperature+1081.08166425371000000000000
            # density = float(f"{densityf:.5g}")# 示例值，需要根据实际情况修改
            # viscosity = float(f"{viscosityf:.5g}")
            # logging.info(f"工质 {workingfluid} 使用预设物性参数：密度={density} kg/m³，粘度={viscosity} Pa·s")
            # 对于其他工质（如50EG），使用默认值或单独的处理逻辑
            # 特殊温度点处理
            if temperature == -35:  # -35℃处理
                viscosity = 0.06693
                density = 1089.94
            elif temperature == -30:  # -30℃处理
                viscosity = 0.04398
                density = 1089.04
            elif temperature == -25:
                viscosity = 0.0305
                density = 1088.01
            elif temperature == -20:
                viscosity = 0.02207
                density = 1086.87
            elif temperature == -15:
                viscosity = 0.01653
                density = 1085.61
            elif temperature == -10:
                viscosity = 0.01274
                density = 1084.22
            elif temperature == -5:
                viscosity = 0.01005
                density = 1082.71
            elif temperature == 0:
                viscosity = 0.00809
                density = 1081.08
            elif temperature == 5:
                viscosity = 0.00663
                density = 1079.33
            elif temperature == 10:
                viscosity = 0.0055
                density = 1077.46
            elif temperature == 15:
                viscosity = 0.00463
                density = 1075.46
            elif temperature == 20:
                viscosity =0.00394
                density =1073.35
            elif temperature == 25:
                viscosity =0.00339
                density =1071.11
            elif temperature == 30:
                viscosity =0.00294
                density =1068.75
            elif temperature == 35:
                viscosity =0.00256
                density =1066.27
            elif temperature == 40:
                viscosity =0.00226
                density =1063.66
            elif temperature == 45:
                viscosity =0.002
                density =1060.94
            elif temperature == 50:
                viscosity =0.00178
                density =1058.09
            elif temperature == 55:
                viscosity =0.00159
                density =1055.13
            elif temperature == 60:
                viscosity =0.00143
                density =1052.04
            elif temperature == 65:
                viscosity =0.00129
                density =1048.83
            elif temperature == 70:
                viscosity =0.00117
                density =1045.49
            elif temperature == 75:
                viscosity =0.00107
                density =1042.04
            elif temperature == 80:
                viscosity =0.00098
                density =1038.46
            elif temperature == 85:
                viscosity =0.00089
                density =1034.77
            elif temperature == 90:
                viscosity =0.00082
                density =1030.95
            elif temperature == 95:
                viscosity =0.00076
                density =1027.01
            elif temperature == 100:
                viscosity =0.0007
                density =1022.95
            elif temperature == 105:
                viscosity =1018.76
                density =0.00065
            elif temperature == 110:
                viscosity =0.0006
                density =1014.46
            elif temperature == 115:
                viscosity =0.00056
                density =1010.03
            elif temperature == 120:
                viscosity =0.00053
                density =1005.48
            elif temperature == 125:
                viscosity =0.00049
                density =1000.81
            else:
                # 粘度多项式公式
                viscosityf = (
                        -0.000000000000002891086856001190294218 * temperature ** 7
                        + 0.000000000001077136958500754662705 * temperature ** 6
                        - 0.000000000151185766170254471016 * temperature ** 5
                        + 0.00000000991277733151393456151 * temperature ** 4
                        - 0.0000003214977807661548076655 * temperature ** 3
                        + 0.000007525808299778965107627 * temperature ** 2
                        - 0.0002892859986891702268885 * temperature
                        + 0.008247019621432158734131
                )
                viscosity = float(f"{viscosityf:.5g}")

                # 密度计算公式
                densityf = (
                        -0.00000000373756542382253 * temperature ** 3
                        - 0.00243168387588982000000 * temperature ** 2
                        - 0.33814727958631500000000 * temperature
                        + 1081.08166425371000000000000
                )
                density = float(f"{densityf:.5g}")

            # 记录计算方式
            calc_method = "预设值" if temperature in [-35, -30, -25, -20, -15, -10, -5, 0, 5, 10, 15, 20, 25, 30, 35, 40, 45, 50, 55, 60, 65, 70, 75, 80, 85, 90, 95, 100, 105, 110, 115, 120, 125] else "多项式公式"
            logging.info(
                f"工质 {workingfluid} 在{temperature}℃使用{calc_method}计算：密度={density} kg/m³，粘度={viscosity} Pa·s")


        name = extract_model_name(model_import_path)
        datenow=get_formatted_date()
        unicode_operator_name = change_unicode(operator_name)
        public_unicode=change_unicode('仿真自动化结果')

        # folder_path_report = "D:\\STARCCM Simulation automation\\Report"
        # # 创建模型专属报告文件夹
        # report_subfolder = os.path.normpath(os.path.join(folder_path_report, name))
        # if not os.path.exists(report_subfolder):
        #     os.makedirs(report_subfolder)
        #     logging.info(f"已创建模型专属报告文件夹: {report_subfolder}")
        # else:
        #     logging.info(f"报告文件夹已存在: {report_subfolder}")

        report_subfolder=f"D:\\STARCCM Simulation automation\\{datenow}\\{operator_name}\\{name}_{self.index}\\Report"
        report_subfolder_public=f"D:\\仿真自动化结果\\{datenow}\\{operator_name}\\{name}_{self.index}\\Report"

        # 创建文件处理器（按模型名生成日志文件）
        log_path = rf"D:\\STARCCM Simulation automation\\{datenow}\\{operator_name}\\{name}_{self.index}\\Log\\{name}.log"
        file_handler = logging.FileHandler(log_path, encoding='utf-8')
        file_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))

        log_path_public = rf"D:\\仿真自动化结果\\{datenow}\\{operator_name}\\{name}_{self.index}\\Log\\{name}.log"
        file_handler_public = logging.FileHandler(log_path_public, encoding='utf-8')
        file_handler_public.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))

        # 添加处理器
        self.logger.addHandler(file_handler)
        self.logger.addHandler(file_handler_public)

        # 在这里处理这些值，例如启动仿真
        logging.info(f'导入数模路径: {model_import_path}')
        logging.info(f'STAR-CCM软件路径: {starccm_path}')
        logging.info(f'REFPRP64.DLL路径: {refprop_path}')
        logging.info(f'线程数: {threads}')
        logging.info(f'停止准则 最大步数: {stop_criteria_max_steps}')
        logging.info(f'基础尺寸: {base_size}')
        logging.info(f'目标表面尺寸 基数百分比: {target_surface_ratio}')
        logging.info(f'最小表面尺寸 基数百分比: {min_surface_ratio}')
        logging.info(f'棱柱层总厚度 基数百分比: {prisma_layer_thickness_ratio}')
        logging.info(f'棱柱层 层数: {prisma_layer_extension}')
        logging.info(f'入口温度（°C）: {temperature}')
        logging.info(f'入口绝对压力（MPa）: {pressure}')
        logging.info(f'入口质量流量（kg/s）: {inlet_mass_flow_rate}')
        logging.info(f'流体工质: {workingfluid}')
        logging.info(f'动力粘度（Pa·s）: {viscosity}')
        logging.info(f'密度（kg/m³）: {density}')

        rename_and_save_step_file(model_import_path, rf"D:/STARCCM Simulation automation/{datenow}/{operator_name}/{name}_{self.index}/CacheModels", "CacheModel.STEP")
        # rf"D:/STARCCM Simulation automation/{datenow}/{operator_name}/{name}/CacheModels"

        # #宏文件路径
        # if getattr(sys, 'frozen', False):
        #     current_dir = getattr(sys, '_MEIPASS', os.path.dirname(sys.executable))
        # else:
        #     current_dir = os.path.dirname(os.path.abspath(__file__))
        current_dir = os.path.dirname(os.path.abspath(__file__))
        script_path = os.path.normpath(os.path.join(current_dir, 'StarCCM_script.java'))
        # 建议添加路径存在性验证：
        os.makedirs(os.path.dirname(script_path), exist_ok=True)
        # 检查文件是否存在
        if os.path.exists(script_path):
            logging.info(f"文件 {script_path} 存在，将被重写。")
        else:
            logging.info(f"文件 {script_path} 不存在，将被创建。")

        # 写入无后缀名的文件，内容为给定的Java代码
        file_content = rf"""// Simcenter STAR-CCM+ macro: StarCCM_script.java
        // Written by Simcenter STAR-CCM+ 16.06.008
        package macro;
        
        import java.util.*;
        
        import star.base.neo.*;
        import star.segregatedflow.*;
        import star.turbulence.*;
        import star.flow.*;
        import star.energy.*;
        import star.metrics.*;
        import star.meshing.*;
        import star.common.*;
        import star.material.*;
        import star.keturb.*;
        import star.base.report.*;
        import star.prismmesher.*;
        import star.vis.*;
        import star.surfacewrapper.*;
        
        public class StarCCM_script extends StarMacro {{
        
          public void execute() {{
            execute0();
            // D:/STARCCM Simulation automation/{datenow}/{unicode_operator_name}/{name}_{self.index}/Simulation/{name}.sim
            execute1();
            // D:/STARCCM Simulation automation/{datenow}/{unicode_operator_name}/{name}_{self.index}/Simulation/{name}.sim
            execute2();
            // D:/STARCCM Simulation automation/{datenow}/{unicode_operator_name}/{name}_{self.index}/Simulation/{name}.sim
            execute3();
            // D:/STARCCM Simulation automation/{datenow}/{unicode_operator_name}/{name}_{self.index}/Simulation/{name}.sim
            execute4();
            // D:/STARCCM Simulation automation/{datenow}/{unicode_operator_name}/{name}_{self.index}/Simulation/{name}.sim
            execute5();
            // D:/STARCCM Simulation automation/{datenow}/{unicode_operator_name}/{name}_{self.index}/Simulation/{name}.sim
            execute6();
            // D:/STARCCM Simulation automation/{datenow}/{unicode_operator_name}/{name}_{self.index}/Simulation/{name}.sim
            execute7();
          }}
        
          private void execute0() {{
        
            Simulation simulation_0 =
              getActiveSimulation();
        
            PartImportManager partImportManager_0 =
              simulation_0.get(PartImportManager.class);
        
            partImportManager_0.importCadPart(resolvePath("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\CacheModels\\CacheModel.STEP"), "SharpEdges", 30.0, 2, true, 1.0E-5, true, false, false, false, true, NeoProperty.fromString("{{\'STEP\': 0, \'NX\': 0, \'CATIAV5\': 0, \'SE\': 0, \'JT\': 0}}"), true, false);
        
            simulation_0.getSceneManager().createGeometryScene("\u51E0\u4F55\u573A\u666F", "\u8F6E\u5ED3", "\u8868\u9762", 1);
        
            Scene scene_0 =
              simulation_0.getSceneManager().getScene("\u51E0\u4F55\u573A\u666F 1");
        
            scene_0.initializeAndWait();
        
            SceneUpdate sceneUpdate_0 =
              scene_0.getSceneUpdate();
        
            HardcopyProperties hardcopyProperties_0 =
              sceneUpdate_0.getHardcopyProperties();
        
            hardcopyProperties_0.setCurrentResolutionWidth(25);
        
            hardcopyProperties_0.setCurrentResolutionHeight(25);
        
            hardcopyProperties_0.setCurrentResolutionWidth(758);
        
            hardcopyProperties_0.setCurrentResolutionHeight(1191);
        
            scene_0.resetCamera();
        
            simulation_0.saveState("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Simulation\\{name}.sim");
          }}
        
          private void execute1() {{
        
            Simulation simulation_0 =
              getActiveSimulation();
        
            CadPart cadPart_0 =
              ((CadPart) simulation_0.get(SimulationPartManager.class).getPart("Fluid"));
        
            SurfaceWrapperAutoMeshOperation surfaceWrapperAutoMeshOperation_0 =
              (SurfaceWrapperAutoMeshOperation) simulation_0.get(MeshOperationManager.class).createSurfaceWrapperAutoMeshOperation(new NeoObjectVector(new Object[] {{cadPart_0}}), "\u5305\u9762");
        
            surfaceWrapperAutoMeshOperation_0.getDefaultValues().get(BaseSize.class).setValue(1.0);
        
            Units units_0 =
              ((Units) simulation_0.getUnitsManager().getObject("mm"));
        
            surfaceWrapperAutoMeshOperation_0.getDefaultValues().get(BaseSize.class).setUnits(units_0);
        
            surfaceWrapperAutoMeshOperation_0.getDefaultValues().get(BaseSize.class).setValue(1.0);
        
            surfaceWrapperAutoMeshOperation_0.getDefaultValues().get(BaseSize.class).setUnits(units_0);
        
            PartsTargetSurfaceSize partsTargetSurfaceSize_0 =
              surfaceWrapperAutoMeshOperation_0.getDefaultValues().get(PartsTargetSurfaceSize.class);
        
            partsTargetSurfaceSize_0.getRelativeSizeScalar().setValue(25.0);
        
            Units units_1 =
              ((Units) simulation_0.getUnitsManager().getObject(""));
        
            partsTargetSurfaceSize_0.getRelativeSizeScalar().setUnits(units_1);
        
            PartsMinimumSurfaceSize partsMinimumSurfaceSize_0 =
              surfaceWrapperAutoMeshOperation_0.getDefaultValues().get(PartsMinimumSurfaceSize.class);
        
            partsMinimumSurfaceSize_0.getRelativeSizeScalar().setValue(15.0);
        
            partsMinimumSurfaceSize_0.getRelativeSizeScalar().setUnits(units_1);
        
            MeshOperationPart meshOperationPart_0 =
              ((MeshOperationPart) simulation_0.get(SimulationPartManager.class).getPart("\u5305\u9762"));
        
            AutoMeshOperation autoMeshOperation_0 =
              simulation_0.get(MeshOperationManager.class).createAutoMeshOperation(new StringVector(new String[] {{"star.resurfacer.ResurfacerAutoMesher", "star.resurfacer.AutomaticSurfaceRepairAutoMesher", "star.dualmesher.DualAutoMesher", "star.prismmesher.PrismAutoMesher"}}), new NeoObjectVector(new Object[] {{meshOperationPart_0}}));
        
            autoMeshOperation_0.getMesherParallelModeOption().setSelected(MesherParallelModeOption.Type.PARALLEL);
        
            autoMeshOperation_0.getDefaultValues().get(BaseSize.class).setValue({base_size});
        
            autoMeshOperation_0.getDefaultValues().get(BaseSize.class).setUnits(units_0);
        
            PartsTargetSurfaceSize partsTargetSurfaceSize_1 =
              autoMeshOperation_0.getDefaultValues().get(PartsTargetSurfaceSize.class);
        
            partsTargetSurfaceSize_1.getRelativeSizeScalar().setValue({target_surface_ratio});
        
            partsTargetSurfaceSize_1.getRelativeSizeScalar().setUnits(units_1);
        
            PartsMinimumSurfaceSize partsMinimumSurfaceSize_1 =
              autoMeshOperation_0.getDefaultValues().get(PartsMinimumSurfaceSize.class);
        
            partsMinimumSurfaceSize_1.getRelativeSizeScalar().setValue({min_surface_ratio});
        
            partsMinimumSurfaceSize_1.getRelativeSizeScalar().setUnits(units_1);
        
            NumPrismLayers numPrismLayers_0 =
              autoMeshOperation_0.getDefaultValues().get(NumPrismLayers.class);
        
            IntegerValue integerValue_0 =
              numPrismLayers_0.getNumLayersValue();
        
            integerValue_0.getQuantity().setValue({prisma_layer_extension});
        
            PrismLayerStretching prismLayerStretching_0 =
              autoMeshOperation_0.getDefaultValues().get(PrismLayerStretching.class);
        
            prismLayerStretching_0.getStretchingQuantity().setValue(1.3);
        
            prismLayerStretching_0.getStretchingQuantity().setUnits(units_1);
        
            PrismThickness prismThickness_0 =
              autoMeshOperation_0.getDefaultValues().get(PrismThickness.class);
        
            prismThickness_0.getRelativeSizeScalar().setValue({prisma_layer_thickness_ratio});
        
            prismThickness_0.getRelativeSizeScalar().setUnits(units_1);
        
            MaximumCellSize maximumCellSize_0 =
              autoMeshOperation_0.getDefaultValues().get(MaximumCellSize.class);
        
            maximumCellSize_0.getRelativeSizeScalar().setValue(100.0);
        
            maximumCellSize_0.getRelativeSizeScalar().setUnits(units_1);
        
            PhysicsContinuum physicsContinuum_0 =
              simulation_0.getContinuumManager().createContinuum(PhysicsContinuum.class);
        
            physicsContinuum_0.enable(ThreeDimensionalModel.class);
        
            physicsContinuum_0.enable(SingleComponentGasModel.class);
        
            physicsContinuum_0.enable(SegregatedFlowModel.class);
        
            physicsContinuum_0.enable(ConstantDensityModel.class);
        
            physicsContinuum_0.enable(SteadyModel.class);
        
            physicsContinuum_0.enable(TurbulentModel.class);
        
            physicsContinuum_0.enable(RansTurbulenceModel.class);
        
            physicsContinuum_0.enable(KEpsilonTurbulence.class);
        
            physicsContinuum_0.enable(RkeTwoLayerTurbModel.class);
        
            physicsContinuum_0.enable(KeTwoLayerAllYplusWallTreatment.class);
        
            SingleComponentGasModel singleComponentGasModel_0 =
              physicsContinuum_0.getModelManager().getModel(SingleComponentGasModel.class);
        
            Gas gas_0 =
              ((Gas) singleComponentGasModel_0.getMaterial());
        
            ConstantMaterialPropertyMethod constantMaterialPropertyMethod_0 =
              ((ConstantMaterialPropertyMethod) gas_0.getMaterialProperties().getMaterialProperty(DynamicViscosityProperty.class).getMethod());
        
            constantMaterialPropertyMethod_0.getQuantity().setValue({viscosity});
        
            Units units_2 =
              ((Units) simulation_0.getUnitsManager().getObject("Pa-s"));
        
            constantMaterialPropertyMethod_0.getQuantity().setUnits(units_2);
        
            ConstantMaterialPropertyMethod constantMaterialPropertyMethod_1 =
              ((ConstantMaterialPropertyMethod) gas_0.getMaterialProperties().getMaterialProperty(ConstantDensityProperty.class).getMethod());
        
            constantMaterialPropertyMethod_1.getQuantity().setValue({density});
        
            Units units_3 =
              ((Units) simulation_0.getUnitsManager().getObject("kg/m^3"));
        
            constantMaterialPropertyMethod_1.getQuantity().setUnits(units_3);
        
            physicsContinuum_0.getReferenceValues().get(ReferencePressure.class).setValue(0.0);
        
            Units units_4 =
              ((Units) simulation_0.getUnitsManager().getObject("Pa"));
        
            physicsContinuum_0.getReferenceValues().get(ReferencePressure.class).setUnits(units_4);
        
            InitialPressureProfile initialPressureProfile_0 =
              physicsContinuum_0.getInitialConditions().get(InitialPressureProfile.class);
        
            initialPressureProfile_0.getMethod(ConstantScalarProfileMethod.class).getQuantity().setValue({pressure});
        
            Units units_5 =
              ((Units) simulation_0.getUnitsManager().getObject("MPa"));
        
            initialPressureProfile_0.getMethod(ConstantScalarProfileMethod.class).getQuantity().setUnits(units_5);
        
            simulation_0.getRegionManager().newRegionsFromParts(new NeoObjectVector(new Object[] {{meshOperationPart_0}}), "OneRegionPerPart", null, "OneBoundaryPerPartSurface", null, "OneFeatureCurve", null, RegionManager.CreateInterfaceMode.BOUNDARY, "OneEdgeBoundaryPerPart", null);
        
            Region region_0 =
              simulation_0.getRegionManager().getRegion("\u5305\u9762");
        
            Boundary boundary_0 =
              region_0.getBoundaryManager().getBoundary("Fluid.inlet");
        
            MassFlowBoundary massFlowBoundary_0 =
              ((MassFlowBoundary) simulation_0.get(ConditionTypeManager.class).get(MassFlowBoundary.class));
        
            boundary_0.setBoundaryType(massFlowBoundary_0);
        
            Boundary boundary_1 =
              region_0.getBoundaryManager().getBoundary("Fluid.outlet");
        
            OutletBoundary outletBoundary_0 =
              ((OutletBoundary) simulation_0.get(ConditionTypeManager.class).get(OutletBoundary.class));
        
            boundary_1.setBoundaryType(outletBoundary_0);
        
            Units units_6 =
              simulation_0.getUnitsManager().getInternalUnits(new IntVector(new int[] {{1, 0, -1, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0}}));
        
            MassFlowRateProfile massFlowRateProfile_0 =
              boundary_0.getValues().get(MassFlowRateProfile.class);
        
            massFlowRateProfile_0.getMethod(ConstantScalarProfileMethod.class).getQuantity().setDefinition("{inlet_mass_flow_rate}");
        
            massFlowRateProfile_0.getMethod(ConstantScalarProfileMethod.class).getQuantity().setUnits(units_6);
        
            StepStoppingCriterion stepStoppingCriterion_0 =
              ((StepStoppingCriterion) simulation_0.getSolverStoppingCriterionManager().getSolverStoppingCriterion("Maximum Steps"));
        
            IntegerValue integerValue_1 =
              stepStoppingCriterion_0.getMaximumNumberStepsObject();
        
            integerValue_1.getQuantity().setValue({stop_criteria_max_steps});
        
            PressureDropReport pressureDropReport_0 =
              simulation_0.getReportManager().createReport(PressureDropReport.class);
        
            pressureDropReport_0.setPresentationName("Dp");
        
            pressureDropReport_0.getParts().setQuery(null);
        
            pressureDropReport_0.getParts().setObjects(boundary_0);
        
            pressureDropReport_0.getLowPressureParts().setQuery(null);
        
            pressureDropReport_0.getLowPressureParts().setObjects(boundary_1);
        
            LatestMeshProxyRepresentation latestMeshProxyRepresentation_0 =
              ((LatestMeshProxyRepresentation) simulation_0.getRepresentationManager().getObject("Latest Surface/Volume"));
        
            pressureDropReport_0.setRepresentation(latestMeshProxyRepresentation_0);
        
            simulation_0.getMonitorManager().createMonitorAndPlot(new NeoObjectVector(new Object[] {{pressureDropReport_0}}), true, "%1$s \u7ED8\u56FE");
        
            ReportMonitor reportMonitor_0 =
              ((ReportMonitor) simulation_0.getMonitorManager().getMonitor("Dp Monitor"));
        
            MonitorPlot monitorPlot_0 =
              simulation_0.getPlotManager().createMonitorPlot(new NeoObjectVector(new Object[] {{reportMonitor_0}}), "Dp Monitor \u7ED8\u56FE");
        
            monitorPlot_0.open();
        
            PlotUpdate plotUpdate_0 =
              monitorPlot_0.getPlotUpdate();
        
            HardcopyProperties hardcopyProperties_1 =
              plotUpdate_0.getHardcopyProperties();
        
            hardcopyProperties_1.setCurrentResolutionWidth(25);
        
            hardcopyProperties_1.setCurrentResolutionHeight(25);
        
            Scene scene_0 =
              simulation_0.getSceneManager().getScene("\u51E0\u4F55\u573A\u666F 1");
        
            SceneUpdate sceneUpdate_0 =
              scene_0.getSceneUpdate();
        
            HardcopyProperties hardcopyProperties_0 =
              sceneUpdate_0.getHardcopyProperties();
        
            hardcopyProperties_0.setCurrentResolutionWidth(760);
        
            hardcopyProperties_0.setCurrentResolutionHeight(1192);
        
            hardcopyProperties_1.setCurrentResolutionWidth(758);
        
            hardcopyProperties_1.setCurrentResolutionHeight(1191);
        
            StatisticsReport statisticsReport_0 =
              simulation_0.getReportManager().createReport(StatisticsReport.class);
        
            statisticsReport_0.setPresentationName("A_dp");
        
            statisticsReport_0.setSampleFilterOption(SampleFilterOption.LastNSamples);
        
            statisticsReport_0.setMonitor(reportMonitor_0);
        
            LastNSamplesFilter lastNSamplesFilter_0 =
              ((LastNSamplesFilter) statisticsReport_0.getSampleFilterManager().getObject("\u6700\u540E N \u4E2A\u6837\u672C"));
        
            lastNSamplesFilter_0.setNSamples(200);
        
            simulation_0.getMonitorManager().createMonitorAndPlot(new NeoObjectVector(new Object[] {{statisticsReport_0}}), true, "%1$s \u7ED8\u56FE");
        
            ReportMonitor reportMonitor_1 =
              ((ReportMonitor) simulation_0.getMonitorManager().getMonitor("A_dp Monitor"));
        
            MonitorPlot monitorPlot_1 =
              simulation_0.getPlotManager().createMonitorPlot(new NeoObjectVector(new Object[] {{reportMonitor_1}}), "A_dp Monitor \u7ED8\u56FE");
        
            monitorPlot_1.open();
        
            PlotUpdate plotUpdate_1 =
              monitorPlot_1.getPlotUpdate();
        
            HardcopyProperties hardcopyProperties_2 =
              plotUpdate_1.getHardcopyProperties();
        
            hardcopyProperties_2.setCurrentResolutionWidth(25);
        
            hardcopyProperties_2.setCurrentResolutionHeight(25);
        
            hardcopyProperties_1.setCurrentResolutionWidth(760);
        
            hardcopyProperties_1.setCurrentResolutionHeight(1192);
        
            hardcopyProperties_2.setCurrentResolutionWidth(758);
        
            hardcopyProperties_2.setCurrentResolutionHeight(1191);
        
            simulation_0.saveState("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Simulation\\{name}.sim");
          }}
        
          private void execute2() {{
        
            Simulation simulation_0 =
              getActiveSimulation();
        
            MeshPipelineController meshPipelineController_0 =
              simulation_0.get(MeshPipelineController.class);
        
            meshPipelineController_0.generateVolumeMesh();
        
            simulation_0.saveState("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Simulation\\{name}.sim");
          }}
        
          private void execute3() {{
        
            Simulation simulation_0 =
              getActiveSimulation();
        
            MaxReport maxReport_0 =
              simulation_0.getReportManager().createReport(MaxReport.class);
        
            maxReport_0.setPresentationName("V_max");
        
            PrimitiveFieldFunction primitiveFieldFunction_0 =
              ((PrimitiveFieldFunction) simulation_0.getFieldFunctionManager().getFunction("Velocity"));
        
            VectorMagnitudeFieldFunction vectorMagnitudeFieldFunction_0 =
              ((VectorMagnitudeFieldFunction) primitiveFieldFunction_0.getMagnitudeFunction());
        
            maxReport_0.setFieldFunction(vectorMagnitudeFieldFunction_0);
        
            maxReport_0.getParts().setQuery(null);
        
            Region region_0 =
              simulation_0.getRegionManager().getRegion("\u5305\u9762");
        
            maxReport_0.getParts().setObjects(region_0);
        
            simulation_0.getMonitorManager().createMonitorAndPlot(new NeoObjectVector(new Object[] {{maxReport_0}}), true, "%1$s \u7ED8\u56FE");
        
            ReportMonitor reportMonitor_2 =
              ((ReportMonitor) simulation_0.getMonitorManager().getMonitor("V_max Monitor"));
        
            MonitorPlot monitorPlot_2 =
              simulation_0.getPlotManager().createMonitorPlot(new NeoObjectVector(new Object[] {{reportMonitor_2}}), "V_max Monitor \u7ED8\u56FE");
        
            monitorPlot_2.open();
        
            PlotUpdate plotUpdate_2 =
              monitorPlot_2.getPlotUpdate();
        
            HardcopyProperties hardcopyProperties_3 =
              plotUpdate_2.getHardcopyProperties();
        
            hardcopyProperties_3.setCurrentResolutionWidth(25);
        
            hardcopyProperties_3.setCurrentResolutionHeight(25);
        
            MonitorPlot monitorPlot_1 =
              ((MonitorPlot) simulation_0.getPlotManager().getPlot("A_dp Monitor \u7ED8\u56FE"));
        
            PlotUpdate plotUpdate_1 =
              monitorPlot_1.getPlotUpdate();
        
            HardcopyProperties hardcopyProperties_2 =
              plotUpdate_1.getHardcopyProperties();
        
            hardcopyProperties_2.setCurrentResolutionWidth(760);
        
            hardcopyProperties_2.setCurrentResolutionHeight(1192);
        
            hardcopyProperties_3.setCurrentResolutionWidth(758);
        
            hardcopyProperties_3.setCurrentResolutionHeight(1191);
        
            simulation_0.saveState("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Simulation\\{name}.sim");
          }}
        
          private void execute4() {{
        
            Simulation simulation_0 =
              getActiveSimulation();
        
            ResidualPlot residualPlot_0 =
              ((ResidualPlot) simulation_0.getPlotManager().getPlot("Residuals"));
        
            residualPlot_0.open();
        
            PlotUpdate plotUpdate_3 =
              residualPlot_0.getPlotUpdate();
        
            HardcopyProperties hardcopyProperties_4 =
              plotUpdate_3.getHardcopyProperties();
        
            hardcopyProperties_4.setCurrentResolutionWidth(25);
        
            hardcopyProperties_4.setCurrentResolutionHeight(25);
        
            simulation_0.getSimulationIterator().run();
        
            MonitorPlot monitorPlot_2 =
              ((MonitorPlot) simulation_0.getPlotManager().getPlot("V_max Monitor \u7ED8\u56FE"));
        
            PlotUpdate plotUpdate_2 =
              monitorPlot_2.getPlotUpdate();
        
            HardcopyProperties hardcopyProperties_3 =
              plotUpdate_2.getHardcopyProperties();
        
            hardcopyProperties_3.setCurrentResolutionWidth(760);
        
            hardcopyProperties_3.setCurrentResolutionHeight(1192);
        
            hardcopyProperties_4.setCurrentResolutionWidth(758);
        
            hardcopyProperties_4.setCurrentResolutionHeight(1191);
        
            MonitorPlot monitorPlot_0 =
              ((MonitorPlot) simulation_0.getPlotManager().getPlot("Dp Monitor \u7ED8\u56FE"));
        
            PlotUpdate plotUpdate_0 =
              monitorPlot_0.getPlotUpdate();
        
            HardcopyProperties hardcopyProperties_1 =
              plotUpdate_0.getHardcopyProperties();
        
            hardcopyProperties_1.setCurrentResolutionWidth(758);
        
            hardcopyProperties_1.setCurrentResolutionHeight(1191);
        
            simulation_0.saveState("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Simulation\\{name}.sim");
          }}
        
          private void execute5() {{
        
            Simulation simulation_0 =
              getActiveSimulation();
        
            MonitorPlot monitorPlot_0 =
              ((MonitorPlot) simulation_0.getPlotManager().getPlot("Dp Monitor \u7ED8\u56FE"));
        
            Cartesian2DAxisManager cartesian2DAxisManager_0 =
              ((Cartesian2DAxisManager) monitorPlot_0.getAxisManager());
        
            cartesian2DAxisManager_0.setAxesBounds(new Vector(Arrays.<AxisManager.AxisBounds>asList(new AxisManager.AxisBounds("Left Axis", -737733.1683996408, false, 102205.11569759721, false), new AxisManager.AxisBounds("Bottom Axis", 1.0, false, {stop_criteria_max_steps}, false))));
        
            monitorPlot_0.export(resolvePath("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_pressure.csv"), ",");
            
            monitorPlot_0.export(resolvePath("D:\\{public_unicode}\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_pressure.csv"), ",");
        
            MonitorPlot monitorPlot_1 =
              ((MonitorPlot) simulation_0.getPlotManager().getPlot("A_dp Monitor \u7ED8\u56FE"));
        
            PlotUpdate plotUpdate_1 =
              monitorPlot_1.getPlotUpdate();
        
            HardcopyProperties hardcopyProperties_2 =
              plotUpdate_1.getHardcopyProperties();
        
            hardcopyProperties_2.setCurrentResolutionWidth(758);
        
            hardcopyProperties_2.setCurrentResolutionHeight(1191);
        
            Cartesian2DAxisManager cartesian2DAxisManager_1 =
              ((Cartesian2DAxisManager) monitorPlot_1.getAxisManager());
        
            cartesian2DAxisManager_1.setAxesBounds(new Vector(Arrays.<AxisManager.AxisBounds>asList(new AxisManager.AxisBounds("Left Axis", -480055.1526588006, false, 31154.726553345095, false), new AxisManager.AxisBounds("Bottom Axis", 1.0, false, {stop_criteria_max_steps}, false))));
        
            monitorPlot_1.export(resolvePath("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_average_pressure.csv"), ",");
            
            monitorPlot_1.export(resolvePath("D:\\{public_unicode}\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_average_pressure.csv"), ",");
        
            MonitorPlot monitorPlot_2 =
              ((MonitorPlot) simulation_0.getPlotManager().getPlot("V_max Monitor \u7ED8\u56FE"));
        
            PlotUpdate plotUpdate_2 =
              monitorPlot_2.getPlotUpdate();
        
            HardcopyProperties hardcopyProperties_3 =
              plotUpdate_2.getHardcopyProperties();
        
            hardcopyProperties_3.setCurrentResolutionWidth(758);
        
            hardcopyProperties_3.setCurrentResolutionHeight(1191);
        
            Cartesian2DAxisManager cartesian2DAxisManager_2 =
              ((Cartesian2DAxisManager) monitorPlot_2.getAxisManager());
        
            cartesian2DAxisManager_2.setAxesBounds(new Vector(Arrays.<AxisManager.AxisBounds>asList(new AxisManager.AxisBounds("Left Axis", 80.78857937401325, false, 1560.7257699953568, false), new AxisManager.AxisBounds("Bottom Axis", 1.0, false, 101.0, false))));
        
            monitorPlot_2.export(resolvePath("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_V_max.csv"), ",");
            
            monitorPlot_2.export(resolvePath("D:\\{public_unicode}\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_V_max.csv"), ",");
        
            simulation_0.getSceneManager().createScalarScene("\u6807\u91CF\u573A\u666F", "\u8F6E\u5ED3", "\u6807\u91CF");
        
            Scene scene_1 =
              simulation_0.getSceneManager().getScene("\u6807\u91CF\u573A\u666F 1");
        
            scene_1.initializeAndWait();
        
            ScalarDisplayer scalarDisplayer_0 =
              ((ScalarDisplayer) scene_1.getDisplayerManager().getObject("\u6807\u91CF 1"));
        
            Legend legend_0 =
              scalarDisplayer_0.getLegend();
        
            PredefinedLookupTable predefinedLookupTable_0 =
              ((PredefinedLookupTable) simulation_0.get(LookupTableManager.class).getObject("blue-yellow-red"));
        
            legend_0.setLookupTable(predefinedLookupTable_0);
        
            SceneUpdate sceneUpdate_1 =
              scene_1.getSceneUpdate();
        
            HardcopyProperties hardcopyProperties_5 =
              sceneUpdate_1.getHardcopyProperties();
        
            hardcopyProperties_5.setCurrentResolutionWidth(25);
        
            hardcopyProperties_5.setCurrentResolutionHeight(25);
        
            hardcopyProperties_3.setCurrentResolutionWidth(760);
        
            hardcopyProperties_3.setCurrentResolutionHeight(1192);
        
            ResidualPlot residualPlot_0 =
              ((ResidualPlot) simulation_0.getPlotManager().getPlot("Residuals"));
        
            PlotUpdate plotUpdate_3 =
              residualPlot_0.getPlotUpdate();
        
            HardcopyProperties hardcopyProperties_4 =
              plotUpdate_3.getHardcopyProperties();
        
            hardcopyProperties_4.setCurrentResolutionWidth(760);
        
            hardcopyProperties_4.setCurrentResolutionHeight(1192);
        
            hardcopyProperties_5.setCurrentResolutionWidth(758);
        
            hardcopyProperties_5.setCurrentResolutionHeight(1191);
        
            scene_1.resetCamera();
        
            scene_1.setPresentationName("\u538B\u529B\u4E91\u56FE");
        
            LogoAnnotation logoAnnotation_0 =
              ((LogoAnnotation) simulation_0.getAnnotationManager().getObject("Logo"));
        
            logoAnnotation_0.setOpacity(0.0);
        
            scalarDisplayer_0.getInputParts().setQuery(null);
        
            Region region_0 =
              simulation_0.getRegionManager().getRegion("\u5305\u9762");
        
            Boundary boundary_2 =
              region_0.getBoundaryManager().getBoundary("Fluid.Faces");
        
            Boundary boundary_0 =
              region_0.getBoundaryManager().getBoundary("Fluid.inlet");
        
            Boundary boundary_1 =
              region_0.getBoundaryManager().getBoundary("Fluid.outlet");
        
            scalarDisplayer_0.getInputParts().setObjects(boundary_2, boundary_0, boundary_1);
        
            PrimitiveFieldFunction primitiveFieldFunction_1 =
              ((PrimitiveFieldFunction) simulation_0.getFieldFunctionManager().getFunction("AbsolutePressure"));
        
            scalarDisplayer_0.getScalarDisplayQuantity().setFieldFunction(primitiveFieldFunction_1);
        
            scalarDisplayer_0.setFillMode(ScalarFillMode.NODE_FILLED);
        
            BlueRedLookupTable blueRedLookupTable_0 =
              ((BlueRedLookupTable) simulation_0.get(LookupTableManager.class).getObject("blue-red"));
        
            legend_0.setLookupTable(blueRedLookupTable_0);
        
            legend_0.setTitleHeight(0.035);
        
            legend_0.setLabelHeight(0.035);
        
            legend_0.setWidth(0.35);
        
            legend_0.setPositionCoordinate(new DoubleVector(new double[] {{0.4, 0.08}}));
        
            legend_0.setLabelFormat({dp_format});
        
            legend_0.setNumberOfLabels(5);
        
            Units units_7 =
              ((Units) simulation_0.getUnitsManager().getObject({dp_unit}));
        
            scalarDisplayer_0.getScalarDisplayQuantity().setUnits(units_7);
        
            CurrentView currentView_0 =
              scene_1.getCurrentView();
        
            currentView_0.setInput(new DoubleVector(new double[] {{0.006064277158636892, 0.005535606369249629, 0.07199999063106509}}), new DoubleVector(new double[] {{0.006064277158636892, 0.005535606369249629, 0.47804024423186187}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}), 0.10599777638115217, 1, 30.0);
        
            scene_1.setViewOrientation(new DoubleVector(new double[] {{1.0, 1.0, 1.0}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}));
        
            scene_1.resetCamera();
        
            //currentView_0.setInput(new DoubleVector(new double[] {{0.006064277158636892, 0.005535606369249629, 0.07199999063106509}}), new DoubleVector(new double[] {{0.24049172687688083, 0.23996305608749355, 0.306427440349309}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}), 0.06688408114046082, 1, 30.0);
        
            scene_1.printAndWait(resolvePath("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_\u538B\u529B\u4E91\u56FE.png"), 2, 1600, 900, true, false);
            
            scene_1.printAndWait(resolvePath("D:\\{public_unicode}\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_\u538B\u529B\u4E91\u56FE.png"), 2, 1600, 900, true, false);
        
            //currentView_0.setInput(new DoubleVector(new double[] {{0.006064277158636892, 0.005535606369249629, 0.07199999063106509}}), new DoubleVector(new double[] {{0.24049172687688083, 0.23996305608749355, 0.306427440349309}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}), 0.06688408114046082, 1, 30.0);
        
            scene_1.export3DSceneFileAndWait(resolvePath("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_\u538B\u529B\u4E91\u56FE.sce"), "\u538B\u529B\u4E91\u56FE", "", false, SceneFileCompressionLevel.OFF);
            
            scene_1.export3DSceneFileAndWait(resolvePath("D:\\{public_unicode}\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_\u538B\u529B\u4E91\u56FE.sce"), "\u538B\u529B\u4E91\u56FE", "", false, SceneFileCompressionLevel.OFF);
        
            Units units_8 =
              simulation_0.getUnitsManager().getPreferredUnits(Dimensions.Builder().length(1).build());
        
            scene_1.setTransparencyOverrideMode(SceneTransparencyOverride.MAKE_SCENE_TRANSPARENT);
        
            scene_1.getCreatorGroup().setQuery(null);
        
            scene_1.getCreatorGroup().setObjects(region_0);
        
            scene_1.getCreatorGroup().setQuery(null);
        
            scene_1.getCreatorGroup().setObjects(region_0);
        
            scene_1.getCreatorGroup().setQuery(null);
        
            scene_1.getCreatorGroup().setObjects(region_0);
        
            PrimitiveFieldFunction primitiveFieldFunction_0 =
              ((PrimitiveFieldFunction) simulation_0.getFieldFunctionManager().getFunction("Velocity"));
        
            StreamPart streamPart_0 =
              simulation_0.getPartManager().createStreamPart(new NeoObjectVector(new Object[] {{region_0}}), new NeoObjectVector(new Object[] {{boundary_2, boundary_0, boundary_1}}), primitiveFieldFunction_0, 8, 8, 2);
        
            scene_1.setTransparencyOverrideMode(SceneTransparencyOverride.USE_DISPLAYER_PROPERTY);
        
            simulation_0.getSceneManager().createEmptyScene("\u573A\u666F");
        
            Scene scene_2 =
              simulation_0.getSceneManager().getScene("\u573A\u666F 1");
        
            scene_2.initializeAndWait();
        
            SceneUpdate sceneUpdate_2 =
              scene_2.getSceneUpdate();
        
            HardcopyProperties hardcopyProperties_6 =
              sceneUpdate_2.getHardcopyProperties();
        
            hardcopyProperties_6.setCurrentResolutionWidth(25);
        
            hardcopyProperties_6.setCurrentResolutionHeight(25);
        
            hardcopyProperties_5.setCurrentResolutionWidth(760);
        
            hardcopyProperties_5.setCurrentResolutionHeight(1192);
        
            hardcopyProperties_6.setCurrentResolutionWidth(758);
        
            hardcopyProperties_6.setCurrentResolutionHeight(1191);
        
            scene_2.resetCamera();
        
            scene_2.setPresentationName("\u6D41\u7EBF\u56FE");
        
            ScalarDisplayer scalarDisplayer_1 =
              scene_2.getDisplayerManager().createScalarDisplayer("\u6807\u91CF");
        
            Legend legend_1 =
              scalarDisplayer_1.getLegend();
        
            legend_1.setLookupTable(predefinedLookupTable_0);
        
            simulation_0.getSceneManager().deleteScenes(new NeoObjectVector(new Object[] {{scene_2}}));
        
            hardcopyProperties_5.setCurrentResolutionWidth(758);
        
            hardcopyProperties_5.setCurrentResolutionHeight(1191);
        
            simulation_0.getSceneManager().createEmptyScene("\u573A\u666F");
        
            Scene scene_3 =
              simulation_0.getSceneManager().getScene("\u573A\u666F 1");
        
            scene_3.initializeAndWait();
        
            SceneUpdate sceneUpdate_3 =
              scene_3.getSceneUpdate();
        
            HardcopyProperties hardcopyProperties_7 =
              sceneUpdate_3.getHardcopyProperties();
        
            hardcopyProperties_7.setCurrentResolutionWidth(25);
        
            hardcopyProperties_7.setCurrentResolutionHeight(25);
        
            hardcopyProperties_5.setCurrentResolutionWidth(760);
        
            hardcopyProperties_5.setCurrentResolutionHeight(1192);
        
            hardcopyProperties_7.setCurrentResolutionWidth(758);
        
            hardcopyProperties_7.setCurrentResolutionHeight(1191);
        
            scene_3.resetCamera();
        
            scene_3.setPresentationName("\u6D41\u7EBF\u56FE");
        
            StreamDisplayer streamDisplayer_0 =
              scene_3.getDisplayerManager().createStreamDisplayer("\u6D41\u7EBF");
        
            Legend legend_2 =
              streamDisplayer_0.getLegend();
        
            legend_2.setLookupTable(predefinedLookupTable_0);
        
            streamDisplayer_0.getInputParts().setQuery(null);
        
            streamDisplayer_0.getInputParts().setObjects(streamPart_0);
        
            VectorMagnitudeFieldFunction vectorMagnitudeFieldFunction_0 =
              ((VectorMagnitudeFieldFunction) primitiveFieldFunction_0.getMagnitudeFunction());
        
            streamDisplayer_0.getScalarDisplayQuantity().setFieldFunction(vectorMagnitudeFieldFunction_0);
        
            streamDisplayer_0.setMode(StreamDisplayerMode.LINES);
        
            legend_2.setLookupTable(blueRedLookupTable_0);
        
            legend_2.setTitleHeight(0.035);
        
            legend_2.setLabelHeight(0.035);
        
            legend_2.setPositionCoordinate(new DoubleVector(new double[] {{0.4, 0.08}}));
        
            legend_2.setLabelFormat("%-6.2f");
        
            legend_2.setNumberOfLabels(5);
        
            legend_2.setWidth(0.35);
        
            PartDisplayer partDisplayer_0 =
              scene_3.getDisplayerManager().createPartDisplayer("\u8868\u9762", -1, 1);
        
            partDisplayer_0.getInputParts().setQuery(null);
        
            FeatureCurve featureCurve_0 =
              ((FeatureCurve) region_0.getFeatureCurveManager().getObject("Default Feature Curve"));
        
            partDisplayer_0.getInputParts().setObjects(boundary_2, boundary_0, boundary_1, featureCurve_0);
        
            partDisplayer_0.setOpacity(0.2);
        
            scene_3.setViewOrientation(new DoubleVector(new double[] {{1.0, 1.0, 1.0}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}));
        
            CurrentView currentView_1 =
              scene_3.getCurrentView();
        
            currentView_1.setInput(new DoubleVector(new double[] {{0.0, 0.0, 0.0}}), new DoubleVector(new double[] {{3.0354027944862283, 3.0354027944862283, 3.0354027944862283}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}), 1.3724755655678504, 1, 30.0);
        
            scene_3.resetCamera();
        
            //currentView_1.setInput(new DoubleVector(new double[] {{0.006064277158636892, 0.005535606369249629, 0.07199999063106509}}), new DoubleVector(new double[] {{0.2404917268768808, 0.23996305608749355, 0.306427440349309}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}), 0.06688408114046082, 1, 30.0);
        
            simulation_0.saveState("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Simulation\\{name}.sim");
          }}
        
          private void execute6() {{
        
            Simulation simulation_0 =
              getActiveSimulation();
        
            Scene scene_3 =
              simulation_0.getSceneManager().getScene("\u6D41\u7EBF\u56FE");
        
            CurrentView currentView_1 =
              scene_3.getCurrentView();
        
            currentView_1.setInput(new DoubleVector(new double[] {{0.006064277158636892, 0.005535606369249629, 0.07199999063106509}}), new DoubleVector(new double[] {{0.2404917268768808, 0.23996305608749355, 0.306427440349309}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}), 0.06688408114046082, 1, 30.0);
        
            scene_3.resetCamera();
        
            scene_3.printAndWait(resolvePath("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_\u6D41\u7EBF\u56FE.png"), 2, 1600, 900, true, false);
            
            scene_3.printAndWait(resolvePath("D:\\{public_unicode}\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_\u6D41\u7EBF\u56FE.png"), 2, 1600, 900, true, false);
        
            //currentView_1.setInput(new DoubleVector(new double[] {{0.006064277158636892, 0.005535606369249629, 0.07199999063106509}}), new DoubleVector(new double[] {{0.2404917268768808, 0.23996305608749355, 0.306427440349309}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}), 0.06688408114046082, 1, 30.0);
        
            scene_3.export3DSceneFileAndWait(resolvePath("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_\u6D41\u7EBF\u56FE.sce"), "\u6D41\u7EBF\u56FE", "", false, SceneFileCompressionLevel.OFF);
            
            scene_3.export3DSceneFileAndWait(resolvePath("D:\\{public_unicode}\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_\u6D41\u7EBF\u56FE.sce"), "\u6D41\u7EBF\u56FE", "", false, SceneFileCompressionLevel.OFF);
        
            Units units_1 =
              simulation_0.getUnitsManager().getPreferredUnits(Dimensions.Builder().build());
        
            scene_3.setTransparencyOverrideMode(SceneTransparencyOverride.MAKE_SCENE_TRANSPARENT);
        
            scene_3.getCreatorGroup().setQuery(null);
        
            Region region_0 =
              simulation_0.getRegionManager().getRegion("\u5305\u9762");
        
            scene_3.getCreatorGroup().setObjects(region_0);
        
            Units units_9 =
              simulation_0.getUnitsManager().getPreferredUnits(Dimensions.Builder().velocity(1).build());
        
            scene_3.getCreatorGroup().setQuery(null);
        
            scene_3.getCreatorGroup().setObjects(region_0);
        
            PrimitiveFieldFunction primitiveFieldFunction_0 =
              ((PrimitiveFieldFunction) simulation_0.getFieldFunctionManager().getFunction("Velocity"));
        
            VectorMagnitudeFieldFunction vectorMagnitudeFieldFunction_0 =
              ((VectorMagnitudeFieldFunction) primitiveFieldFunction_0.getMagnitudeFunction());
        
            ThresholdPart thresholdPart_0 =
              simulation_0.getPartManager().createThresholdPart(new NeoObjectVector(new Object[] {{region_0}}), new DoubleVector(new double[] {{{speed_of_sound*0.3}, 300.0}}), units_9, vectorMagnitudeFieldFunction_0, 0);
        
            scene_3.setTransparencyOverrideMode(SceneTransparencyOverride.USE_DISPLAYER_PROPERTY);
        
            thresholdPart_0.setPresentationName("Ma>0.3");
        
            thresholdPart_0.setPresentationName("Ma>0.3\u533A\u57DF");
        
            simulation_0.getSceneManager().createEmptyScene("\u573A\u666F");
        
            Scene scene_4 =
              simulation_0.getSceneManager().getScene("\u573A\u666F 1");
        
            scene_4.initializeAndWait();
        
            SceneUpdate sceneUpdate_4 =
              scene_4.getSceneUpdate();
        
            HardcopyProperties hardcopyProperties_8 =
              sceneUpdate_4.getHardcopyProperties();
        
            hardcopyProperties_8.setCurrentResolutionWidth(25);
        
            hardcopyProperties_8.setCurrentResolutionHeight(25);
        
            SceneUpdate sceneUpdate_3 =
              scene_3.getSceneUpdate();
        
            HardcopyProperties hardcopyProperties_7 =
              sceneUpdate_3.getHardcopyProperties();
        
            hardcopyProperties_7.setCurrentResolutionWidth(760);
        
            hardcopyProperties_7.setCurrentResolutionHeight(1192);
        
            hardcopyProperties_8.setCurrentResolutionWidth(758);
        
            hardcopyProperties_8.setCurrentResolutionHeight(1191);
        
            scene_4.resetCamera();
        
            scene_4.setPresentationName("Ma>0.3\u533A\u57DF\u56FE");
        
            PartDisplayer partDisplayer_1 =
              scene_4.getDisplayerManager().createPartDisplayer("\u8868\u9762", -1, 1);
        
            partDisplayer_1.getInputParts().setQuery(null);
        
            Boundary boundary_2 =
              region_0.getBoundaryManager().getBoundary("Fluid.Faces");
        
            Boundary boundary_0 =
              region_0.getBoundaryManager().getBoundary("Fluid.inlet");
        
            Boundary boundary_1 =
              region_0.getBoundaryManager().getBoundary("Fluid.outlet");
        
            FeatureCurve featureCurve_0 =
              ((FeatureCurve) region_0.getFeatureCurveManager().getObject("Default Feature Curve"));
        
            partDisplayer_1.getInputParts().setObjects(boundary_2, boundary_0, boundary_1, featureCurve_0);
        
            partDisplayer_1.setOpacity(0.2);
        
            PartDisplayer partDisplayer_2 =
              scene_4.getDisplayerManager().createPartDisplayer("\u8868\u9762", -1, 1);
        
            partDisplayer_2.getInputParts().setQuery(null);
        
            partDisplayer_2.getInputParts().setObjects(thresholdPart_0);
        
            CurrentView currentView_2 =
              scene_4.getCurrentView();
        
            currentView_2.setInput(new DoubleVector(new double[] {{0.0, 0.0, 0.0}}), new DoubleVector(new double[] {{0.0, 0.0, 5.257471861486698}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}), 1.3724755655678502, 1, 30.0);
        
            scene_4.setViewOrientation(new DoubleVector(new double[] {{1.0, 1.0, 1.0}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}));
        
            scene_4.resetCamera();
        
            //currentView_2.setInput(new DoubleVector(new double[] {{0.006064277158636892, 0.005535606369249629, 0.07199999063106509}}), new DoubleVector(new double[] {{0.2404917268768808, 0.23996305608749355, 0.306427440349309}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}), 0.06688408114046082, 1, 30.0);
        
            scene_4.printAndWait(resolvePath("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_Ma_0.3\u533A\u57DF\u56FE.png"), 2, 1600, 900, true, false);
            
            scene_4.printAndWait(resolvePath("D:\\{public_unicode}\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_Ma_0.3\u533A\u57DF\u56FE.png"), 2, 1600, 900, true, false);
        
            //currentView_2.setInput(new DoubleVector(new double[] {{0.006064277158636892, 0.005535606369249629, 0.07199999063106509}}), new DoubleVector(new double[] {{0.2404917268768808, 0.23996305608749355, 0.306427440349309}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}), 0.06688408114046082, 1, 30.0);
        
            scene_4.export3DSceneFileAndWait(resolvePath("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_Ma_0.3\u533A\u57DF\u56FE.sce"), "Ma>0.3\u533A\u57DF\u56FE", "", false, SceneFileCompressionLevel.OFF);
            
            scene_4.export3DSceneFileAndWait(resolvePath("D:\\{public_unicode}\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_Ma_0.3\u533A\u57DF\u56FE.sce"), "Ma>0.3\u533A\u57DF\u56FE", "", false, SceneFileCompressionLevel.OFF);
        
            simulation_0.saveState("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Simulation\\{name}.sim");
          }}
        
          private void execute7() {{
            Simulation simulation_0 =
              getActiveSimulation();

            simulation_0.getSceneManager().createGeometryScene("\u51E0\u4F55\u573A\u666F", "\u8F6E\u5ED3", "\u8868\u9762", 1);

            Scene scene_5 =
              simulation_0.getSceneManager().getScene("\u51E0\u4F55\u573A\u666F 2");

            scene_5.initializeAndWait();

            SceneUpdate sceneUpdate_2 =
              scene_5.getSceneUpdate();

            HardcopyProperties hardcopyProperties_2 =
              sceneUpdate_2.getHardcopyProperties();

            hardcopyProperties_2.setCurrentResolutionWidth(25);

            hardcopyProperties_2.setCurrentResolutionHeight(25);

            Scene scene_0 =
              simulation_0.getSceneManager().getScene("\u51E0\u4F55\u573A\u666F 1");

            SceneUpdate sceneUpdate_0 =
              scene_0.getSceneUpdate();

            HardcopyProperties hardcopyProperties_0 =
              sceneUpdate_0.getHardcopyProperties();

            hardcopyProperties_0.setCurrentResolutionWidth(1818);

            hardcopyProperties_0.setCurrentResolutionHeight(856);

            hardcopyProperties_2.setCurrentResolutionWidth(1816);

            hardcopyProperties_2.setCurrentResolutionHeight(855);

            scene_5.resetCamera();

            scene_5.setPresentationName("\u6D41\u4F53\u57DF\u56FE");

            CurrentView currentView_3 =
              scene_5.getCurrentView();

            //currentView_3.setInput(new DoubleVector(new double[] {{0.0060642761908053285, 0.005535606360364112, 0.07199999063106509}}), new DoubleVector(new double[] {{0.0060642761908053285, 0.005535606360364112, 0.3304198118061474}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}), 0.06746111050434028, 1, 30.0);

            scene_5.setViewOrientation(new DoubleVector(new double[] {{1.0, 1.0, 1.0}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}));

            scene_5.setTransparencyOverrideMode(SceneTransparencyOverride.MAKE_SCENE_TRANSPARENT);

            scene_5.resetCamera();

            //currentView_3.setInput(new DoubleVector(new double[] {{0.0060642761908053285, 0.005535606360364112, 0.07199999063106509}}), new DoubleVector(new double[] {{0.15526302951017404, 0.1547343596797328, 0.22119874395043382}}), new DoubleVector(new double[] {{0.0, 1.0, 0.0}}), 0.06688397135209896, 1, 30.0);

            scene_5.printAndWait(resolvePath("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_\u6D41\u4F53\u57DF\u56FE.png"), 2, 1600, 900, true, false);
            
            scene_5.printAndWait(resolvePath("D:\\{public_unicode}\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Report\\{name}_\u6D41\u4F53\u57DF\u56FE.png"), 2, 1600, 900, true, false);
            
            simulation_0.saveState("D:\\STARCCM Simulation automation\\{datenow}\\{unicode_operator_name}\\{name}_{self.index}\\Simulation\\{name}.sim");
          }}
        }}"""

        # 逐行写入文件
        with open(script_path, "w", encoding="utf-8") as file:
            for line in file_content.splitlines():
                file.write(line + "\n")

        # 构建命令
        command = [
            starccm_path,
            "-verbose", # 强制输出详细日志
            "-np", f"{threads}",  # 使用指定的处理器核心数
            "-batch",  # 批处理模式
            # "-new",
            # "-macro",
            script_path
        ]

        # # 执行命令
        # result = subprocess.run(command, shell=True)

        # 修改原有的subprocess.run调用方式
        process = subprocess.Popen(
            command,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            universal_newlines=True,
            bufsize=1,
            encoding='utf-8',
            errors='replace'
        )
        # 实时捕获输出
        while True:
            output = process.stdout.readline()
            if output == '' and process.poll() is not None:
                break
            if output:
                # 同时输出到控制台和文件
                self.logger.info(output.strip())  # 替换原来的logging.info
                QApplication.processEvents()  # 保持UI响应

        # 获取最终返回码
        result = process.poll()

        # 检查命令执行结果
        if process.returncode == 0:

            self.process_state = 1

            end_time = datetime.datetime.now()
            duration = end_time - start_time
            duration = f"{duration.seconds // 3600}小时{(duration.seconds // 60) % 60}分{duration.seconds % 60}秒"

            # # 仅非队列任务更新参数记录
            # if not is_queue_task:
            #     self.last_input_params = current_params.copy()
            #     self.save_config()

            # 创建图片显示区域
            image_container = QWidget()
            image_layout = QHBoxLayout(image_container)

            # 压力云图
            pressure_img = os.path.join(report_subfolder_public, f"{name}_压力云图.png")
            if os.path.exists(pressure_img):
                pixmap = QPixmap(pressure_img).scaled(480, 270, Qt.KeepAspectRatio)
                self.pressure_label.setPixmap(pixmap)
                self.pressure_label.mouseDoubleClickEvent = lambda e: self.show_image(pressure_img)

            # 流线图
            streamline_img = os.path.join(report_subfolder_public, f"{name}_流线图.png")
            if os.path.exists(streamline_img):
                pixmap = QPixmap(streamline_img).scaled(480, 270, Qt.KeepAspectRatio)
                self.streamline_label.setPixmap(pixmap)
                self.streamline_label.mouseDoubleClickEvent = lambda e: self.show_image(streamline_img)

            model_img=os.path.join(report_subfolder_public, f"{name}_流体域图.png")
            Ma_img = os.path.join(report_subfolder_public, f"{name}_Ma_0.3区域图.png")

            # 将图片容器添加到主界面
            self.layout().addWidget(image_container)

            logging.info("仿真成功完成")
            logging.info(f"报告已保存至 D:\\STARCCM Simulation automation\\{datenow}\\{operator_name}\\{name}_{self.index}\\Report\\{name}_pressure.csv")
            logging.info(f"报告已保存至 D:\\仿真自动化结果\\{datenow}\\{operator_name}\\{name}_{self.index}\\Report\\{name}_pressure.csv")

            # 仿真完成后恢复按钮状态
            self.run_button.setText('开始运行')
            self.run_button.setStyleSheet("font-size: 48px; font-weight: bold; background-color: #90EE90;")  # 恢复绿色背景
            self.run_button.setEnabled(True)
            QApplication.processEvents()  # 强制刷新UI

            Ma = 0  # 是否删除Ma>0.3图，0删除

            #最大流速读取
            vmax_file_path = f"D:\\STARCCM Simulation automation\\{datenow}\\{operator_name}\\{name}_{self.index}\\Report\\{name}_V_max.csv"
            vmax = read_last_row_last_column(vmax_file_path)
            if workingfluid in ["R134a", "R1234yf","R744"]:
                vmax_float = float(vmax)
                # 分级判断流速范围
                if vmax_float > speed_of_sound * 0.5:
                    Ma = 2
                    logging.warning(
                        f"{name}内部最大流速为: {vmax_float:.3g} m/s (Ma={vmax_float / speed_of_sound:.2f}), 内部部分流速马赫数超0.5！")
                    logging.info(
                        f"Ma>0.3区域图已保存至 D:\\STARCCM Simulation automation\\{datenow}\\{operator_name}\\{name}_{self.index}\\Report\\{name}_Ma_0.3区域图.png")
                    logging.info(
                        f"Ma>0.3区域图已保存至 D:\\仿真自动化结果\\{datenow}\\{operator_name}\\{name}_{self.index}\\Report\\{name}_Ma_0.3区域图.png")
                elif vmax_float > speed_of_sound * 0.3:
                    Ma = 1
                    logging.warning(
                        f"{name}内部最大流速为: {vmax_float:.3g} m/s (Ma={vmax_float / speed_of_sound:.2f}), 内部部分流速马赫数超0.3")
                    logging.info(
                        f"Ma>0.3区域图已保存至 D:\\STARCCM Simulation automation\\{datenow}\\{operator_name}\\{name}_{self.index}\\Report\\{name}_Ma_0.3区域图.png")
                    logging.info(
                        f"Ma>0.3区域图已保存至 D:\\仿真自动化结果\\{datenow}\\{operator_name}\\{name}_{self.index}\\Report\\{name}_Ma_0.3区域图.png")
                else:
                    Ma = 0
                    logging.info(
                        f"{name}内部最大流速为: {vmax_float:.3g} m/s (Ma={vmax_float / speed_of_sound:.2f}), 内部流速马赫数未超过0.3")
            else:
                Ma = 0

            if workingfluid in ["R134a", "R1234yf","R744"]:
                res_mach_number=round(float(vmax) / speed_of_sound, 3)
            else:
                res_mach_number="N/A"
            if is_queue_task:
                # 如果是队列任务，找到对应的任务字典
                pending_tasks = [t for t in self.task_queue if t['status'] == "计算中"]
                if pending_tasks:
                    task = pending_tasks[0]
                    # 添加index和date到任务字典
                    task['simulation_index'] = self.index
                    task['simulation_date'] = datenow
                    task['Ma'] = Ma
                    task['res_mach_number'] = res_mach_number
                    # self.save_config()  # 立即保存配置

                # 在成功运行后更新参数记录
            self.last_input_params = current_params.copy()
            # 保存输入条件
            self.save_config()

            self.res_sce_path1 = self.get_scene_path("压力云图.sce",current_params)
            self.res_sce_path2= self.get_scene_path("流线图.sce",current_params)
            self.res_sce_path3= self.get_scene_path("Ma_0.3区域图.sce",current_params)

            self.btn_mach_3d.setVisible(Ma != 0)
            self.btn_pressure_3d.setEnabled(True)
            self.btn_streamline_3d.setEnabled(True)
            self.btn_mach_3d.setEnabled(True if Ma != 0 else False)

            # 刷新界面确保按钮状态更新
            self.btn_pressure_3d.repaint()
            self.btn_streamline_3d.repaint()
            self.btn_mach_3d.repaint()

            # 计算结果并格式化
            calculated_value = safe_eval(inlet_mass_flow_rate)
            formatted_value = f"{calculated_value:.2f}" if calculated_value is not None else "N/A"
            csv_file_path = f"D:\\STARCCM Simulation automation\\{datenow}\\{operator_name}\\{name}_{self.index}\\Report\\{name}_average_pressure.csv"
            last_value = read_last_row_last_column(csv_file_path)
            output_pptname=f'{datenow}_{operator_name}_{name}_{self.index}.pptx'
            output_pptpath=os.path.join(report_subfolder,output_pptname)
            output_pptpath_public = os.path.join(report_subfolder_public,output_pptname)

            def resource_path(relative_path):
                if hasattr(sys, '_MEIPASS'):
                    return os.path.join(sys._MEIPASS, relative_path)
                return os.path.join(os.path.abspath("."), relative_path)

                # 修改模板加载方式

            #自动输出PPT部分
            if workingfluid in ["R134a","R1234yf","R744"]:
                ppt_template = 'Refrigerant_Report.pptx'
                prs = Presentation(str(resource_path(ppt_template)))
                slide_1 = prs.slides[0]
                slide_2 = prs.slides[1]
                slide_3 = prs.slides[2]
                slide_4 = prs.slides[3]
                slide_5 = prs.slides[4]
                append_text_to_slide(slide_1,"压降仿真_",f'{datenow}_{operator_name}_{name}_{self.index}')
                modify_table(slide_3,2,1, f'{workingfluid}',0)
                modify_table(slide_3,2,2, f'{temperature}',0)
                modify_table(slide_3,2,3, f'{pressure}',0)
                modify_table(slide_3,2,4, f'{density}',0)
                modify_table(slide_3,2,5, f'{viscosity}',0)
                replace_image(slide_3, "",model_img)
                modify_table(slide_4,2,1, f'{self.index}',0)
                modify_table(slide_4,2,2, f'{pressure}',0)
                modify_table(slide_4,2,3, f'{temperature}',0)
                modify_table(slide_4,2,4, f'{formatted_value}',0)
                modify_table(slide_4,2,5,  f"{float(last_value)/1000:.2f}",0)
                replace_image(slide_4, "",pressure_img,0)#替换第一张
                replace_image(slide_4, "",streamline_img,0)#由于第一张被替换，第二张变成了第一张，所以再次替换
                replace_image(slide_5, "",Ma_img)
                if Ma==0:
                    del slide_5
            else:
                ppt_template = '50EG_Report.pptx'
                prs = Presentation(str(resource_path(ppt_template)))
                slide_1 = prs.slides[0]
                slide_2 = prs.slides[1]
                slide_3 = prs.slides[2]
                slide_4 = prs.slides[3]
                append_text_to_slide(slide_1, "压降仿真_", f'{datenow}_{operator_name}_{name}_{self.index}')
                modify_table(slide_3, 2, 1, f'{workingfluid}', 0)
                modify_table(slide_3, 2, 2, f'{temperature}', 0)
                modify_table(slide_3, 2, 3, f'{density}', 0)
                modify_table(slide_3, 2, 4, f'{viscosity}', 0)
                replace_image(slide_3, "", model_img)
                modify_table(slide_4, 2, 1, f'{self.index}', 0)
                modify_table(slide_4, 2, 2, f'{temperature}', 0)
                modify_table(slide_4, 2, 3, f'{formatted_value}', 0)
                modify_table(slide_4, 2, 4, f"{float(last_value) / 1000:.2f}", 0)
                replace_image(slide_4, "", pressure_img, 0)
                replace_image(slide_4, "", streamline_img, 0)
            prs.save(output_pptpath)
            prs.save(output_pptpath_public)

            if Ma==0:
                # 定义要删除的文件路径
                ma_image = os.path.join(report_subfolder, f"{name}_Ma_0.3区域图.png")
                ma_sce = os.path.join(report_subfolder, f"{name}_Ma_0.3区域图.sce")
                ma_image_public = os.path.join(report_subfolder_public, f"{name}_Ma_0.3区域图.png")
                ma_sce_public = os.path.join(report_subfolder_public, f"{name}_Ma_0.3区域图.sce")

                # 安全删除文件
                for file_path in [ma_image, ma_sce, ma_image_public, ma_sce_public]:
                    try:
                        if os.path.exists(file_path):
                            os.remove(file_path)
                            # logging.info(f"已删除文件: {file_path}")
                    except Exception as e:
                        logging.error(f"删除文件失败: {file_path}, 错误: {str(e)}")

            if Ma==2:
                dp_image = os.path.join(report_subfolder, f"{name}_压力云图.png")
                dp_sce = os.path.join(report_subfolder, f"{name}_压力云图.sce")
                vline_image = os.path.join(report_subfolder, f"{name}_流线图.png")
                vline_sce = os.path.join(report_subfolder, f"{name}_流线图.sce")
                dp_csv = os.path.join(report_subfolder, f"{name}_pressure.csv")
                adp_csv = os.path.join(report_subfolder, f"{name}_average_pressure.csv")
                dp_image_public = os.path.join(report_subfolder_public, f"{name}_压力云图.png")
                dp_sce_public = os.path.join(report_subfolder_public, f"{name}_压力云图.sce")
                vline_image_public = os.path.join(report_subfolder_public, f"{name}_流线图.png")
                vline_sce_public = os.path.join(report_subfolder_public, f"{name}_流线图.sce")
                dp_csv_public = os.path.join(report_subfolder_public, f"{name}_pressure.csv")
                adp_csv_public = os.path.join(report_subfolder_public, f"{name}_average_pressure.csv")
                for file_path in [dp_image, dp_sce, vline_image, vline_sce, dp_csv, adp_csv, dp_image_public,dp_sce_public, vline_image_public, vline_sce_public, dp_csv_public, adp_csv_public]:
                    try:
                        if os.path.exists(file_path):
                            os.remove(file_path)
                            logging.info(f"已删除文件: {file_path}")
                    except Exception as e:
                        logging.error(f"删除文件失败: {file_path}, 错误: {str(e)}")

            #马赫数获取
            if workingfluid in ["R134a", "R1234yf","R744"] and speed_of_sound > 0:
                try:
                    # 计算马赫数并保留3位有效数字
                    mach_number = round(float(vmax) / speed_of_sound, 3)

                    # 设置显示颜色和文本
                    if mach_number > 0.5:
                        color = "#FF0000"  # 红色
                    elif 0.3 <= mach_number <= 0.5:
                        color = "#FFA500"  # 橙色
                    else:
                        color = "#009900"  # 绿色
                    self.mach_number_label.setText(f"{mach_number}")
                    self.mach_number_label.setStyleSheet(f"""
                        font-size: 24px; 
                        color: {color}; 
                        font-weight: bold;
                    """)
                    logging.info(f"马赫数计算完成：{mach_number}")
                except Exception as e:
                    logging.error(f"马赫数计算失败: {str(e)}")
                    self.mach_number_label.setText("计算错误")
                    self.mach_number_label.setStyleSheet("font-size: 24px; color: #FF0000; font-weight: bold;")
            else:
                # 非压缩性流体隐藏马赫数显示
                self.mach_number_label.setText("N/A")
                self.mach_number_label.setStyleSheet("font-size: 24px; color: #666666; font-weight: bold;")

            # 确保在界面更新后执行
            QApplication.processEvents()

            csv_file_path = f"D:\\STARCCM Simulation automation\\{datenow}\\{operator_name}\\{name}_{self.index}\\Report\\{name}_average_pressure.csv"
            last_value = read_last_row_last_column(csv_file_path)
            if last_value is not None:
                # 更新压降显示
                drop_value = int(round(float(last_value)))
                self.pressure_drop_label.setText(f"{name}压降值为: {drop_value} Pa")
                self.pressure_drop_label.setStyleSheet("font-size: 24px; color: #009900; font-weight: bold;")
                QApplication.processEvents()
                logging.info(f'导入数模路径: {model_import_path}')
                logging.info(f'STAR-CCM软件路径: {starccm_path}')
                logging.info(f'REFPRP64.DLL路径: {refprop_path}')
                logging.info(f'线程数: {threads}')
                logging.info(f'停止准则 最大步数: {stop_criteria_max_steps}')
                logging.info(f'基础尺寸: {base_size}')
                logging.info(f'目标表面尺寸 基数百分比: {target_surface_ratio}')
                logging.info(f'最小表面尺寸 基数百分比: {min_surface_ratio}')
                logging.info(f'棱柱层总厚度 基数百分比: {prisma_layer_thickness_ratio}')
                logging.info(f'棱柱层 层数: {prisma_layer_extension}')
                logging.info(f'入口温度（°C）: {temperature}')
                logging.info(f'入口绝对压力（MPa）: {pressure}')
                logging.info(f'入口质量流量（kg/s）: {inlet_mass_flow_rate}')
                logging.info(f'流体工质: {workingfluid}')
                logging.info(f'动力粘度（Pa·s）: {viscosity}')
                logging.info(f'密度（kg/m³）: {density}')
                logging.info(f"{name}产品压降值为（Pa）: {int(round(float(last_value)))}")
                logging.info(f'PPT报告已保存完成：{output_pptpath_public}')
            else:
                self.pressure_drop_label.setText("N/A 马赫数>0.5")
                self.pressure_drop_label.setStyleSheet("font-size: 24px; color: #FF0000; font-weight: bold;")
                QApplication.processEvents()
                logging.info(f'导入数模路径: {model_import_path}')
                logging.info(f'入口温度（°C）: {temperature}')
                logging.info(f'入口绝对压力（MPa）: {pressure}')
                logging.info(f'入口质量流量（kg/s）: {inlet_mass_flow_rate}')
                logging.info(f'流体工质: {workingfluid}')
                logging.info(f'动力粘度（Pa·s）: {viscosity}')
                logging.info(f'密度（kg/m³）: {density}')
                logging.warning(f'{name}压降值不显示,因为内部部分流速马赫数>0.5')
                logging.info(f'PPT报告已保存完成：{output_pptpath_public}')

            # 生成报告文件
            report_content = f"""仿真报告
========================

基本信息
--------
导入数模路径: {model_import_path}
数模名称: {name}
线程数: {threads}
最大步数: {stop_criteria_max_steps}

流体参数
--------
流体工质: {workingfluid}
入口温度: {temperature} ℃
入口绝对压力: {pressure} MPa
入口质量流量: {inlet_mass_flow_rate} kg/s
动力粘度: {viscosity} Pa·s
密度: {density} kg/m³
声速: {'N/A' if workingfluid == '50EG' else speed_of_sound} m/s

操作信息
--------
操作员姓名: {operator_name}
仿真开始时间: {start_time.strftime('%Y-%m-%d %H:%M:%S')}
仿真结束时间: {end_time.strftime('%Y-%m-%d %H:%M:%S')} 
仿真用时: {duration}

计算结果
--------
最大马赫数: {'N/A' if workingfluid == '50EG' else round(float(vmax)/speed_of_sound,3)}
{self.get_pressure_drop_display(Ma, last_value)}
"""

            report_path = os.path.join(report_subfolder, f"{name}_仿真报告.txt")
            with open(report_path, 'w', encoding='utf-8') as f:
                f.write(report_content)
            # logging.info(f"仿真报告已生成: {report_path}")

            report_path_public= os.path.join(report_subfolder_public, f"{name}_仿真报告.txt")
            with open(report_path_public, 'w', encoding='utf-8') as f:
                f.write(report_content)
            logging.info(f"仿真报告已生成: {report_path_public}")



        else:
            self.process_state=0
            logging.info(f'导入数模路径: {model_import_path}')
            logging.info(f'STAR-CCM软件路径: {starccm_path}')
            logging.info(f'REFPRP64.DLL路径: {refprop_path}')
            logging.info(f'线程数: {threads}')
            logging.info(f'停止准则 最大步数: {stop_criteria_max_steps}')
            logging.info(f'基础尺寸: {base_size}')
            logging.info(f'目标表面尺寸 基数百分比: {target_surface_ratio}')
            logging.info(f'最小表面尺寸 基数百分比: {min_surface_ratio}')
            logging.info(f'棱柱层总厚度 基数百分比: {prisma_layer_thickness_ratio}')
            logging.info(f'棱柱层 层数: {prisma_layer_extension}')
            logging.info(f'入口温度（°C）: {temperature}')
            logging.info(f'入口绝对压力（MPa）: {pressure}')
            logging.info(f'入口质量流量（kg/s）: {inlet_mass_flow_rate}')
            logging.info(f'流体工质: {workingfluid}')
            logging.info(f'动力粘度（Pa·s）: {viscosity}')
            logging.info(f'密度（kg/m³）: {density}')
            logging.error(f"仿真失败，返回码: {process.returncode}")
            self.pressure_drop_label.setText("压降值获取失败")
            self.pressure_drop_label.setStyleSheet("font-size: 24px; color: #FF0000; font-weight: bold;")
            self.mach_number_label.setText("马赫数获取失败")
            self.mach_number_label.setStyleSheet("font-size: 24px; color: #FF0000; font-weight: bold;")
            # 仿真失败后恢复按钮状态
            self.run_button.setText('开始运行')
            self.run_button.setStyleSheet("font-size: 48px; font-weight: bold; background-color: #90EE90;")  # 恢复绿色背景
            self.run_button.setEnabled(True)
            QApplication.processEvents()  # 强制刷新UI
            if is_queue_task:
                # 需要获取当前正在处理的任务对象
                for task in self.task_queue:
                    if task.get('status') == "计算中":
                        task['status'] = "失败"
                        break  # 只更新第一个处于计算中状态的任务





if __name__ == '__main__':
    # app = QApplication(sys.argv)
    # window = SimulationConfigWindow()
    # sys.exit(app.exec_())
    pass